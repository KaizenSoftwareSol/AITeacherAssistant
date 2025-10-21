# services/document_parser.py

import io
import json
import re
from datetime import datetime
from typing import Any, Dict, Optional

import aiohttp
import PyPDF2
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from logger import logger
from models.document import DocumentType, WebsiteContent


class DocumentParser:
    """Service for parsing various document types and extracting content."""

    @staticmethod
    async def parse_document(
        file_content: bytes, document_type: DocumentType, filename: str
    ) -> Dict[str, Any]:
        """
        Parse document content based on type and return structured data.

        Args:
            file_content: Raw file content as bytes
            document_type: Type of document to parse
            filename: Original filename

        Returns:
            Dictionary containing parsed content and metadata
        """
        try:
            # Handle both string and enum
            if isinstance(document_type, str):
                doc_type_str = document_type.lower()
            else:
                doc_type_str = document_type.value.lower()

            logger.info(f"📄 Parsing {doc_type_str} document: {filename}")

            if doc_type_str == "pdf" or document_type == DocumentType.PDF:
                result = await DocumentParser._parse_pdf(file_content, filename)
                logger.info(
                    f"✅ PDF parsed: {result.get('summary', {}).get('total_pages', 0)} pages"
                )
                return result
            elif doc_type_str == "pptx" or document_type == DocumentType.PPTX:
                result = await DocumentParser._parse_pptx(file_content, filename)
                logger.info(
                    f"✅ PPTX parsed: {result.get('summary', {}).get('total_pages', 0)} slides"
                )
                return result
            elif doc_type_str == "docx" or document_type == DocumentType.DOCX:
                result = await DocumentParser._parse_docx(file_content, filename)
                logger.info(
                    f"✅ DOCX parsed: {result.get('summary', {}).get('total_headings', 0)} headings"
                )
                return result
            else:
                raise ValueError(f"Unsupported document type: {document_type}")

        except Exception as e:
            logger.error(f"❌ Error parsing document {filename}: {str(e)}")
            raise Exception(f"Failed to parse document: {str(e)}")

    @staticmethod
    async def _parse_pdf(file_content: bytes, filename: str) -> Dict[str, Any]:
        """Parse PDF content with chapter and section-based organization."""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            # Extract PDF metadata
            pdf_metadata = pdf_reader.metadata if pdf_reader.metadata else {}
            metadata = {
                "filename": filename,
                "file_type": "PDF",
                "parsed_at": datetime.utcnow().isoformat(),
                "total_pages": len(pdf_reader.pages),
                "title": (
                    pdf_metadata.get("/Title", "") if pdf_metadata.get("/Title") else ""
                ),
                "author": (
                    pdf_metadata.get("/Author", "")
                    if pdf_metadata.get("/Author")
                    else ""
                ),
                "subject": (
                    pdf_metadata.get("/Subject", "")
                    if pdf_metadata.get("/Subject")
                    else ""
                ),
                "creator": (
                    pdf_metadata.get("/Creator", "")
                    if pdf_metadata.get("/Creator")
                    else ""
                ),
                "producer": (
                    pdf_metadata.get("/Producer", "")
                    if pdf_metadata.get("/Producer")
                    else ""
                ),
                "creation_date": (
                    str(pdf_metadata.get("/CreationDate", ""))
                    if pdf_metadata.get("/CreationDate")
                    else ""
                ),
                "modification_date": (
                    str(pdf_metadata.get("/ModDate", ""))
                    if pdf_metadata.get("/ModDate")
                    else ""
                ),
            }

            # STEP 1: Extract text from all pages with metadata
            logger.info("📄 Extracting text from all pages...")
            pages_data = []  # Store page data with metadata
            total_words = 0
            pages_with_content = 0
            pages_with_errors = 0

            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        # Clean text to remove surrogate and problematic characters
                        text = text.encode("utf-8", errors="ignore").decode("utf-8")
                        lines = [
                            line.strip() for line in text.split("\n") if line.strip()
                        ]

                        # Check if page is from appendix or exercise solutions
                        is_appendix = DocumentParser._is_appendix_page(lines, page_num)

                        # Detect chapter from footer/header
                        chapter_num = DocumentParser._detect_chapter_from_footer(
                            lines, page_num
                        )

                        pages_data.append(
                            {
                                "page_num": page_num,
                                "lines": lines,
                                "is_appendix": is_appendix,
                                "chapter_num": chapter_num,
                                "word_count": len(text.split()),
                            }
                        )

                        if not is_appendix:
                            total_words += len(text.split())
                            pages_with_content += 1
                except Exception as e:
                    logger.warning(
                        f"Error extracting text from page {page_num}: {str(e)}"
                    )
                    pages_with_errors += 1

            logger.info(
                f"✅ Extracted {len(pages_data)} pages, "
                f"{sum(1 for p in pages_data if not p['is_appendix'])} valid pages"
            )

            # STEP 2: Parse Table of Contents WITH PAGE RANGES
            logger.info("📋 Parsing Table of Contents...")
            # Get all lines from first 50 pages for TOC
            toc_lines = []
            for page_data in pages_data[:50]:
                if not page_data["is_appendix"]:
                    toc_lines.extend(page_data["lines"])

            toc_structure, chapter_ranges = DocumentParser._parse_toc(toc_lines)
            logger.info(f"✅ TOC Method: Found {len(toc_structure)} chapters")

            # FALLBACK: If TOC parsing failed, detect chapters from content
            if not toc_structure:
                logger.warning("TOC parsing failed, using content-based detection...")
                toc_structure, chapter_ranges = (
                    DocumentParser._detect_chapters_from_content(pages_data)
                )
                logger.info(f"✅ Content Method: Found {len(toc_structure)} chapters")

            # STEP 3: Parse content using chapter ranges
            logger.info("📖 Parsing chapter content by page ranges...")
            content_by_chapters = DocumentParser._parse_content_by_pages(
                pages_data, toc_structure, chapter_ranges
            )

            # Compile final result
            result = {
                "type": "pdf",
                "metadata": metadata,
                "content": content_by_chapters,
                "summary": {
                    "total_pages": len(pdf_reader.pages),
                    "pages_with_content": pages_with_content,
                    "pages_with_errors": pages_with_errors,
                    "empty_pages": (
                        len(pdf_reader.pages) - pages_with_content - pages_with_errors
                    ),
                    "total_chapters": len(content_by_chapters),
                    "total_sections": sum(
                        len(sections) if isinstance(sections, dict) else 0
                        for sections in content_by_chapters.values()
                    ),
                },
                "total_word_count": total_words,
            }

            return result

        except Exception as e:
            logger.error(f"Error parsing PDF {filename}: {str(e)}")
            raise Exception(f"Failed to parse PDF: {str(e)}")

    @staticmethod
    def _is_appendix_page(lines: list, page_num: int) -> bool:
        """
        Check if a page is from appendix, exercise solutions, or other back matter.
        Looks at headers, footers, and page content.
        """
        # Check first 10 and last 10 lines
        check_lines = []
        if len(lines) >= 10:
            check_lines.extend(lines[:10])
            check_lines.extend(lines[-10:])
        else:
            check_lines = lines

        # Patterns that indicate appendix/back matter
        appendix_indicators = [
            r"^\s*Appendix\s+[A-Z]",
            r"^\s*APPENDIX\s+[A-Z]",
            r"Exercise\s+Solutions",
            r"EXERCISE\s+SOLUTIONS",
            r"Solutions\s+to\s+Exercise",
            r"Answers\s+to\s+Exercise",
            r"^\s*Index\s*$",
            r"^\s*Glossary\s*$",
            r"^\s*Bibliography\s*$",
            r"^\s*References\s*$",
            r"Appendix\s+[A-Z]:\s+Exercise",
        ]

        for line in check_lines:
            line_clean = line.strip()
            for pattern in appendix_indicators:
                if re.search(pattern, line_clean, re.IGNORECASE):
                    logger.debug(f"Page {page_num} identified as appendix/back matter")
                    return True

        return False

    @staticmethod
    def _detect_chapter_from_footer(lines: list, page_num: int) -> int:
        """
        Detect which chapter a page belongs to by examining footer AND header.
        Uses STRICT patterns to avoid false positives from content references.
        Returns chapter number or None.
        """
        if len(lines) < 2:
            return None

        # ONLY check first 2 and last 2 lines (strict header/footer only)
        header_lines = lines[:2]
        footer_lines = lines[-2:]

        # STRICT patterns - typical book footer/header formats
        # Format: "Page# | Chapter X: Title" or "Chapter X: Title | Page#"
        strict_patterns = [
            r"^\d+\s*\|\s*Chapter\s+(\d+)[:\s]",  # "42 | Chapter 1: ..."
            r"^Chapter\s+(\d+)[:\s].+\|\s*\d+\s*$",  # "Chapter 1: ... | 42"
            r"^\d+\s+\|\s+Chapter\s+(\d+)[:\s]",  # "42  | Chapter 1: ..."
        ]

        # Check ONLY header and footer lines with strict patterns
        for line in header_lines + footer_lines:
            line_clean = line.strip()

            # Must be short (actual footers/headers are concise)
            if len(line_clean) > 120 or len(line_clean) < 10:
                continue

            # Try strict patterns only
            for pattern in strict_patterns:
                match = re.search(pattern, line_clean, re.IGNORECASE)
                if match:
                    chapter_num = int(match.group(1))
                    logger.debug(
                        f"Page {page_num} -> Chapter {chapter_num} "
                        f"[{line_clean[:60]}]"
                    )
                    return chapter_num

        return None

    @staticmethod
    def _parse_toc(lines: list) -> tuple:
        """
        Parse the Table of Contents to extract chapter structure WITH PAGE NUMBERS.
        Returns (toc_dict, chapter_ranges) where:
        - toc_dict: {chapter_num: title}
        - chapter_ranges: {chapter_num: (start_page, end_page)}
        """
        toc_structure = {}
        chapter_pages = {}  # Track start pages from TOC
        in_toc = False

        toc_start_patterns = ["table of contents", "contents"]
        stop_markers = ["appendix", "index", "glossary", "references", "bibliography"]

        # Multiple TOC patterns - MUST handle ALL edge cases
        # Examples: "1 Law's roots 1", "4 Courts 85", "10 'Title' 311"
        toc_patterns = [
            # Explicit "Chapter": "Chapter 1: Title ... 25"
            re.compile(r"^Chapter\s+(\d+)[:\s\-]+(.+?)[\.\s]{2,}.*?(\d+)\s*$"),
            # With dots: "1. Title ... 25"
            re.compile(r"^(\d+)\.\s*(.+?)[\.\s]{2,}.*?(\d+)\s*$"),
            # Simple format with short titles: "1 Law's roots 1", "4 Courts 85"
            # GREEDY match for title, stops at last space + digits
            re.compile(r"^(\d{1,2})\s+(.+?)\s+(\d{1,4})\s*$"),
            # Long spacing: "10     Title     311"
            re.compile(r"^(\d{1,2})\s{2,}(.+?)\s{2,}(\d{1,})\s*$"),
        ]

        for i, line in enumerate(lines):
            line_clean = line.strip()
            line_lower = line_clean.lower()

            # Detect TOC start
            if not in_toc:
                if any(marker in line_lower for marker in toc_start_patterns):
                    if len(line_clean) < 50:
                        in_toc = True
                        logger.debug(f"TOC start at line {i}")
                        continue

            # Process TOC entries
            if in_toc:
                # Stop at appendices
                if any(marker in line_lower for marker in stop_markers):
                    logger.debug(f"TOC end at appendix, line {i}")
                    break

                # Stop at actual Chapter 1 content
                if (
                    line_clean.startswith("Chapter 1")
                    and "..." not in line_clean
                    and not re.search(r"\d+\s*$", line_clean)
                ):
                    logger.debug(f"TOC end at Chapter 1 start")
                    break

                # Try all TOC patterns
                match = None
                for pattern in toc_patterns:
                    match = pattern.match(line_clean)
                    if match:
                        break

                if match:
                    chapter_num = int(match.group(1))
                    chapter_title = match.group(2).strip()
                    page_num = int(match.group(3))

                    # Clean title - remove quotes, extra spaces
                    chapter_title = chapter_title.strip("'\"")
                    chapter_title = re.sub(r"\s{2,}", " ", chapter_title)
                    chapter_title = re.sub(r"[\.\s]+$", "", chapter_title)

                    # Validate chapter number (1-30 reasonable range)
                    if chapter_num < 1 or chapter_num > 30:
                        continue

                    # Validate page number is reasonable (1-2000)
                    if page_num < 1 or page_num > 2000:
                        continue

                    # Validate title - must start with capital and be meaningful
                    if not chapter_title or len(chapter_title) < 2:
                        continue

                    if not chapter_title[0].isupper() and not chapter_title[0] == "'":
                        continue

                    # Skip appendix, exercise, references, etc.
                    skip_words = [
                        "appendix",
                        "exercise",
                        "answer",
                        "reference",
                        "further reading",
                        "index",
                    ]
                    if any(word in chapter_title.lower() for word in skip_words):
                        continue

                    # Skip if already found (avoid duplicates)
                    if chapter_num in toc_structure:
                        continue

                    toc_structure[chapter_num] = chapter_title
                    chapter_pages[chapter_num] = page_num
                    logger.debug(
                        f"TOC: Chapter {chapter_num} '{chapter_title}' -> page {page_num}"
                    )

        # Build chapter ranges (start_page to next_chapter_start - 1)
        chapter_ranges = {}
        sorted_chapters = sorted(chapter_pages.items())

        for i, (chapter_num, start_page) in enumerate(sorted_chapters):
            # End page is start of next chapter - 1 (or max if last chapter)
            if i + 1 < len(sorted_chapters):
                end_page = sorted_chapters[i + 1][1] - 1
            else:
                end_page = 9999  # Will be capped by actual page count

            chapter_ranges[chapter_num] = (start_page, end_page)
            logger.debug(f"Chapter {chapter_num} range: pages {start_page}-{end_page}")

        logger.info(f"Parsed {len(toc_structure)} chapters from TOC with page ranges")
        return toc_structure, chapter_ranges

    @staticmethod
    def _detect_chapters_from_content(pages_data: list) -> tuple:
        """
        Fallback method: Detect chapters by scanning actual content pages.
        STRICT validation to avoid false positives (like years: 1664, 1789).
        Returns (toc_dict, chapter_ranges)
        """
        toc_structure = {}
        chapter_starts = {}

        # FLEXIBLE patterns for content-based detection
        chapter_patterns = [
            (
                r"^Chapter\s+(\d+)[:\s\-]+(.+?)$",
                "explicit_chapter",
            ),  # "Chapter 1: Title"
            (
                r"^CHAPTER\s+(\d+)[:\s\-]+(.+?)$",
                "explicit_chapter",
            ),  # "CHAPTER 1: Title"
            (
                r"^(\d{1,2})\.1\s+(.+)$",
                "first_section",
            ),  # "1.1 Title" = Chapter 1 start
            (
                r"^(\d{1,2})\s*\n+([A-Z][A-Za-z\s]{5,})$",
                "simple_number",
            ),  # "5\nLaw and property"
            (
                r"^(\d{1,2})[\.\s]+([A-Z][^\.]{8,})$",
                "numbered",
            ),  # "5 Law and property" or "5. Title"
            (r"^([IVX]{1,5})[\.\s]+([A-Z][A-Za-z\s]{8,})$", "roman"),  # Roman numerals
        ]

        for page_data in pages_data:
            if page_data["is_appendix"]:
                continue

            page_num = page_data["page_num"]
            lines = page_data["lines"]

            # Skip very early pages (first 10) to avoid TOC
            # But don't skip too much - some books start early
            if page_num < 10:
                continue

            # Check first 10 lines (chapter headings at page top)
            for idx, line in enumerate(lines[:10]):
                line_clean = line.strip()

                # Skip short lines
                if len(line_clean) < 8:
                    continue

                for pattern, pattern_type in chapter_patterns:
                    match = re.match(pattern, line_clean, re.IGNORECASE)
                    if match:
                        chapter_num_raw = match.group(1)
                        chapter_title = (
                            match.group(2).strip() if len(match.groups()) >= 2 else ""
                        )

                        # Convert to integer
                        try:
                            if pattern_type == "roman":
                                # Extended Roman numeral support
                                roman_map = {
                                    "I": 1,
                                    "II": 2,
                                    "III": 3,
                                    "IV": 4,
                                    "V": 5,
                                    "VI": 6,
                                    "VII": 7,
                                    "VIII": 8,
                                    "IX": 9,
                                    "X": 10,
                                    "XI": 11,
                                    "XII": 12,
                                    "XIII": 13,
                                    "XIV": 14,
                                    "XV": 15,
                                    "XVI": 16,
                                    "XVII": 17,
                                    "XVIII": 18,
                                    "XIX": 19,
                                    "XX": 20,
                                }
                                chapter_num = roman_map.get(chapter_num_raw.upper())
                                if chapter_num is None:
                                    continue
                            else:
                                chapter_num = int(chapter_num_raw)
                        except (ValueError, KeyError):
                            continue

                        # **VALIDATION**: Reject unreasonable chapter numbers
                        # Chapters should be 1-30 typically, NOT years like 1664!
                        if chapter_num < 1 or chapter_num > 30:
                            continue

                        # Clean title
                        chapter_title = re.sub(r"[\.\s]+\d+\s*$", "", chapter_title)
                        chapter_title = re.sub(r"\s{2,}", " ", chapter_title).strip()

                        # Validate title - must be meaningful
                        if len(chapter_title) < 3:
                            continue

                        # Skip if already found this chapter
                        if chapter_num in chapter_starts:
                            continue

                        # For first_section pattern, use generic chapter title
                        # because "1.1 Title" will also be a section
                        if pattern_type == "first_section":
                            # Use just the topic from 1.1 as chapter title
                            toc_structure[chapter_num] = chapter_title
                        else:
                            toc_structure[chapter_num] = chapter_title

                        chapter_starts[chapter_num] = page_num
                        logger.info(
                            f"Detected Chapter {chapter_num}: '{chapter_title}' "
                            f"at page {page_num} [{pattern_type}]"
                        )
                        break

        # Build chapter ranges
        chapter_ranges = {}
        sorted_chapters = sorted(chapter_starts.items())

        for i, (chapter_num, start_page) in enumerate(sorted_chapters):
            if i + 1 < len(sorted_chapters):
                end_page = sorted_chapters[i + 1][1] - 1
            else:
                # Last chapter goes to end of non-appendix pages
                non_appendix_pages = [
                    p["page_num"] for p in pages_data if not p["is_appendix"]
                ]
                end_page = (
                    max(non_appendix_pages) if non_appendix_pages else start_page + 10
                )

            chapter_ranges[chapter_num] = (start_page, end_page)
            logger.debug(f"Chapter {chapter_num} range: pages {start_page}-{end_page}")

        return toc_structure, chapter_ranges

    @staticmethod
    def _parse_content_by_pages(
        pages_data: list, toc_structure: dict, chapter_ranges: dict
    ) -> dict:
        """
        Parse content using TOC page ranges - the most reliable method!
        If TOC says "Chapter 1 ... page 35", we know Chapter 1 is on pages 35+
        """
        from collections import OrderedDict

        content_by_chapters = OrderedDict()

        # Sort chapters by number for chronological order
        sorted_chapters = sorted(toc_structure.items())

        # For each chapter, collect pages within its TOC page range
        for chapter_num, chapter_title in sorted_chapters:
            chapter_key = f"Chapter {chapter_num}: {chapter_title}"
            chapter_content = []
            chapter_sections = {}
            current_section = None
            section_buffer = []

            # Get page range from TOC
            start_page, end_page = chapter_ranges.get(chapter_num, (0, 0))

            logger.info(
                f"Parsing Chapter {chapter_num} (pages {start_page}-{end_page})..."
            )

            pages_found = 0

            # Collect all pages in this chapter's range
            for page_data in pages_data:
                page_num = page_data["page_num"]

                # CRITICAL: Skip if appendix page
                if page_data["is_appendix"]:
                    continue

                # Include page if it's in the TOC page range for this chapter
                if start_page <= page_num <= end_page:
                    pages_found += 1
                    # Process lines from this page
                    for line in page_data["lines"]:
                        line_clean = line.strip()

                        # Skip empty lines
                        if not line_clean:
                            continue

                        # Skip chapter heading line itself
                        if re.match(
                            rf"^Chapter\s+{chapter_num}[:\s\-]",
                            line_clean,
                            re.IGNORECASE,
                        ):
                            continue

                        # Skip page numbers and very short lines
                        if len(line_clean) <= 3 or line_clean.isdigit():
                            continue

                        # Skip lines that are just numbers/dots
                        if re.match(r"^[\d\.\s]+$", line_clean):
                            continue

                        # Skip header/footer patterns (page numbers with separators)
                        if re.match(r"^\d+\s*[\|\-]\s*Chapter", line_clean):
                            continue

                        # Skip non-section patterns that look like "X.Y"
                        # FIGURE 1.8, TABLE 2.3, etc. are NOT sections!
                        if re.match(
                            r"^(FIGURE|TABLE|MAP|CHART|IMAGE|FIG\.)\s+\d+\.\d+",
                            line_clean,
                            re.IGNORECASE,
                        ):
                            continue

                        # Check if it's a section heading (X.Y format)
                        section_match = re.match(
                            r"^(\d{1,2})\.(\d{1,2}(?:\.\d+)?)\s+(.+)$", line_clean
                        )

                        if section_match:
                            sec_chapter_num = int(section_match.group(1))
                            sec_num = section_match.group(2)
                            section_title_raw = section_match.group(3).strip()

                            # Validate it belongs to this chapter
                            if sec_chapter_num == chapter_num:
                                # Clean section title - remove bullets and page numbers
                                section_title = re.sub(
                                    r"^[•\-\*]\s*", "", section_title_raw
                                )
                                section_title = re.sub(
                                    r"[\.\s]+\d+\s*$", "", section_title
                                )
                                section_title = re.sub(
                                    r"\s{2,}", " ", section_title
                                ).strip()

                                # CRITICAL: Skip if title starts with common non-section keywords
                                title_lower = section_title.lower()
                                skip_keywords = [
                                    "figure",
                                    "table",
                                    "map",
                                    "chart",
                                    "image",
                                    "the ",
                                    "a ",
                                    "an ",
                                    "this ",
                                    "these ",
                                ]
                                if any(
                                    title_lower.startswith(kw)
                                    for kw in skip_keywords[:5]
                                ):
                                    # Starts with FIGURE, TABLE, etc.
                                    continue

                                # Additional validation - section title should be reasonable
                                if (
                                    len(section_title) > 3
                                    and not section_title.endswith(":")
                                    and "..." not in section_title
                                ):
                                    # Save previous section
                                    if current_section and section_buffer:
                                        chapter_sections[current_section] = " ".join(
                                            section_buffer
                                        ).strip()
                                        section_buffer = []

                                    # Use full section number: "11.3 Title"
                                    full_section_num = f"{sec_chapter_num}.{sec_num}"
                                    current_section = (
                                        f"{full_section_num} {section_title}"
                                    )
                                    logger.debug(f"  Section: {current_section}")
                                    continue

                        # Filter out TOC-like content and noise
                        # Skip lines that look like TOC entries
                        if re.match(r"^\d{1,2}\s+[A-Z].{5,}\s+\d{1,4}\s*$", line_clean):
                            # Looks like "5 Chapter Title 110" = TOC entry
                            continue

                        # Skip lines that are mostly dots (TOC leaders)
                        if line_clean.count(".") > len(line_clean) / 3:
                            continue

                        # Skip very short lines (likely page numbers or headers)
                        if len(line_clean) < 20:
                            continue

                        # Add to content buffer
                        if current_section:
                            section_buffer.append(line_clean)
                        else:
                            chapter_content.append(line_clean)

            # Save last section
            if current_section and section_buffer:
                chapter_sections[current_section] = " ".join(section_buffer).strip()

            logger.info(f"  Found {pages_found} pages for Chapter {chapter_num}")

            # Sort sections by their number (e.g., 1.1, 1.2, 1.3)
            if chapter_sections:

                def get_section_num(sec_key):
                    match = re.match(r"^(\d+)\.(\d+)", sec_key)
                    if match:
                        return (int(match.group(1)), int(match.group(2)))
                    return (999, 999)  # Put non-numbered sections at end

                chapter_sections = dict(
                    sorted(
                        chapter_sections.items(), key=lambda x: get_section_num(x[0])
                    )
                )

            # Add chapter to result
            if chapter_sections:
                content_by_chapters[chapter_key] = chapter_sections
            elif chapter_content:
                content_by_chapters[chapter_key] = {
                    "content": " ".join(chapter_content).strip()
                }
            else:
                logger.warning(
                    f"No content found for {chapter_key} "
                    f"(found {pages_found} pages)"
                )

        # POST-PROCESSING: Remove duplicate sections
        content_by_chapters = DocumentParser._remove_duplicate_sections(
            content_by_chapters
        )

        return dict(content_by_chapters)

    @staticmethod
    def _remove_duplicate_sections(content_by_chapters: dict) -> dict:
        """
        Remove duplicate sections within each chapter.
        Keeps version without bullet points (•, -, *) or the first one found.
        """
        cleaned_chapters = {}

        for chapter_key, sections in content_by_chapters.items():
            if not isinstance(sections, dict):
                cleaned_chapters[chapter_key] = sections
                continue

            # Group sections by their number (e.g., "1.1")
            section_groups = {}
            for section_key, content in sections.items():
                # Extract section number (e.g., "1.1" from "1.1 Title")
                section_num_match = re.match(r"^(\d+\.\d+)", section_key)
                if section_num_match:
                    section_num = section_num_match.group(1)
                    if section_num not in section_groups:
                        section_groups[section_num] = []
                    section_groups[section_num].append(
                        {"key": section_key, "content": content}
                    )
                else:
                    # No number found, keep as is
                    if "other" not in section_groups:
                        section_groups["other"] = []
                    section_groups["other"].append(
                        {"key": section_key, "content": content}
                    )

            # For each section number, pick the best version
            cleaned_sections = {}
            for section_num, duplicates in section_groups.items():
                if len(duplicates) == 1:
                    # No duplicates
                    cleaned_sections[duplicates[0]["key"]] = duplicates[0]["content"]
                else:
                    # Multiple versions - prefer one without bullets
                    best = duplicates[0]
                    for dup in duplicates:
                        # Prefer section without bullet points in the KEY
                        if not re.search(r"[•\-\*]", dup["key"]):
                            best = dup
                            break

                    logger.debug(
                        f"Removed {len(duplicates)-1} duplicate(s) for "
                        f"section {section_num}, kept: {best['key'][:50]}"
                    )
                    cleaned_sections[best["key"]] = best["content"]

            cleaned_chapters[chapter_key] = cleaned_sections

        return cleaned_chapters

    @staticmethod
    def _parse_content_with_toc(
        lines: list,
        toc_structure: dict,
        page_to_line_map: dict,
        page_to_chapter_map: dict,
    ) -> dict:
        """
        Parse document content using TOC structure and footer info as guides.
        Returns content organized by chapters in chronological order.
        Uses page footers to validate chapter assignment.
        """
        from collections import OrderedDict

        content_by_chapters = OrderedDict()

        # Skip TOC and front matter - find where Chapter 1 actually starts
        content_start_idx = 0
        for i, line in enumerate(lines):
            if re.match(r"^Chapter\s+1[:\s\-]", line, re.IGNORECASE):
                # Make sure it's not a TOC entry
                if "..." not in line and not re.search(r"\d{2,}\s*$", line):
                    content_start_idx = i
                    logger.debug(f"Content starts at line {i}")
                    break

        if content_start_idx == 0:
            logger.warning("Could not find Chapter 1 start, using full text")

        # Parse content after TOC
        lines_to_parse = lines[content_start_idx:]

        # Build chapter pattern from TOC
        if toc_structure:
            # Sort chapters by number
            sorted_chapters = sorted(toc_structure.items())
            logger.info(f"Parsing {len(sorted_chapters)} chapters from TOC")
        else:
            logger.warning("No TOC structure found, using pattern matching")
            sorted_chapters = []

        current_chapter = None
        current_chapter_num = None
        current_section = None
        content_buffer = []

        # Section pattern
        section_pattern = re.compile(r"^(\d+\.\d+(?:\.\d+)?)\s+(.+)$")

        # Stop markers for content parsing
        appendix_markers = [
            r"^Appendix\s+[A-Z]",
            r"^APPENDIX\s+[A-Z]",
            r"^Index\s*$",
            r"^References\s*$",
            r"^Bibliography\s*$",
        ]

        for line_idx, line in enumerate(lines_to_parse, start=content_start_idx):
            line_clean = line.strip()

            # Stop parsing at appendices or back matter
            if any(
                re.match(pattern, line_clean, re.IGNORECASE)
                for pattern in appendix_markers
            ):
                logger.info(f"Stopping at appendix/back matter: {line_clean[:50]}")
                break

            # Skip empty or very short lines
            if len(line_clean) <= 3 or line_clean.isdigit():
                continue

            is_chapter = False
            is_section = False

            # Check for chapter headings using TOC structure
            for chapter_num, expected_title in sorted_chapters:
                # Try exact match first
                chapter_pattern = re.compile(
                    rf"^Chapter\s+{chapter_num}[:\s\-]+(.+?)$", re.IGNORECASE
                )
                match = chapter_pattern.match(line_clean)

                if match:
                    found_title = match.group(1).strip()
                    # Clean the found title
                    found_title = re.sub(r"[\.\s]+\d+\s*$", "", found_title)
                    found_title = re.sub(r"\s{2,}", " ", found_title)

                    # Check if it matches TOC (at least partially)
                    title_words = expected_title.lower().split()[:3]
                    found_words = found_title.lower().split()[:3]

                    # If first few words match, it's probably the right chapter
                    if (
                        any(word in found_words for word in title_words)
                        or len(found_title) > 5
                    ):
                        # Save previous chapter content
                        if current_chapter and content_buffer:
                            if current_section:
                                if current_chapter not in content_by_chapters:
                                    content_by_chapters[current_chapter] = {}
                                content_by_chapters[current_chapter][
                                    current_section
                                ] = " ".join(content_buffer).strip()
                            else:
                                if current_chapter not in content_by_chapters:
                                    content_by_chapters[current_chapter] = {}
                                content_by_chapters[current_chapter]["content"] = (
                                    " ".join(content_buffer).strip()
                                )
                            content_buffer = []

                        # Use TOC title (more reliable)
                        current_chapter = f"Chapter {chapter_num}: {expected_title}"
                        current_chapter_num = chapter_num
                        current_section = None
                        is_chapter = True
                        logger.debug(f"Found chapter: {current_chapter}")
                        break

            # Check for section headings (only within a chapter)
            if not is_chapter and current_chapter:
                match = section_pattern.match(line_clean)
                if match:
                    section_num = match.group(1)
                    section_title = match.group(2).strip()

                    # Validate section belongs to current chapter
                    section_chapter_num = int(section_num.split(".")[0])
                    if section_chapter_num == current_chapter_num:
                        # Clean section title
                        section_title = re.sub(r"[\.\s]+\d+\s*$", "", section_title)
                        section_title = re.sub(r"\s{2,}", " ", section_title)

                        # Validate section
                        if (
                            len(section_title) > 3
                            and not section_title.endswith(":")
                            and "..." not in section_title
                        ):
                            # Save previous section
                            if current_section and content_buffer:
                                if current_chapter not in content_by_chapters:
                                    content_by_chapters[current_chapter] = {}
                                content_by_chapters[current_chapter][
                                    current_section
                                ] = " ".join(content_buffer).strip()
                                content_buffer = []

                            current_section = f"{section_num} {section_title}"
                            is_section = True

            # Add to content buffer with footer validation
            if not is_chapter and not is_section and current_chapter:
                # Get the page this line is on
                page_num = page_to_line_map.get(line_idx)

                # Check if page footer confirms we're in the right chapter
                footer_chapter_num = page_to_chapter_map.get(page_num)

                # Only add content if:
                # 1. Line passes basic filters
                # 2. Footer confirms current chapter OR no footer info available
                if (
                    len(line_clean) > 10
                    and not line_clean.isdigit()
                    and not re.match(r"^[\d\.\s]+$", line_clean)
                ):
                    # Validate with footer if available
                    if footer_chapter_num is not None:
                        if footer_chapter_num == current_chapter_num:
                            content_buffer.append(line_clean)
                        else:
                            # Wrong chapter according to footer - skip
                            logger.debug(
                                f"Skipping line from page {page_num}: "
                                f"footer says Ch{footer_chapter_num}, "
                                f"currently parsing Ch{current_chapter_num}"
                            )
                    else:
                        # No footer info, trust our pattern matching
                        content_buffer.append(line_clean)

        # Save final content
        if current_chapter and content_buffer:
            if current_section:
                if current_chapter not in content_by_chapters:
                    content_by_chapters[current_chapter] = {}
                content_by_chapters[current_chapter][current_section] = " ".join(
                    content_buffer
                ).strip()
            else:
                if current_chapter not in content_by_chapters:
                    content_by_chapters[current_chapter] = {}
                content_by_chapters[current_chapter]["content"] = " ".join(
                    content_buffer
                ).strip()

        # If no chapters found, fallback
        if not content_by_chapters:
            logger.warning("No chapters parsed, creating fallback structure")
            content_by_chapters["Document Content"] = {
                "Full Text": " ".join(lines_to_parse)
            }

        return dict(content_by_chapters)  # Convert OrderedDict to dict

    @staticmethod
    async def _parse_pptx(file_content: bytes, filename: str) -> Dict[str, Any]:
        """Parse PowerPoint presentation content with comprehensive extraction."""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)

            # Extract presentation metadata
            metadata = {
                "filename": filename,
                "file_type": "PPTX",
                "parsed_at": datetime.utcnow().isoformat(),
                "total_slides": len(presentation.slides),
                "title": presentation.core_properties.title or "",
                "author": presentation.core_properties.author or "",
                "subject": presentation.core_properties.subject or "",
                "created": (
                    presentation.core_properties.created.isoformat()
                    if presentation.core_properties.created
                    else ""
                ),
                "modified": (
                    presentation.core_properties.modified.isoformat()
                    if presentation.core_properties.modified
                    else ""
                ),
            }

            # Organize content by pages
            pages_content = {}

            for slide_num, slide in enumerate(presentation.slides, 1):
                page_key = f"Page {slide_num}"
                page_content = {
                    "slide_number": slide_num,
                    "title": "",
                    "text_content": [],
                    "images": [],
                    "tables": [],
                    "notes": "",
                    "raw_content": [],
                }

                # Extract slide title (usually from first text box)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not page_content["title"]:
                            page_content["title"] = shape.text.strip()
                        break

                # Extract all content from shapes
                for shape in slide.shapes:
                    shape_data = DocumentParser._extract_shape_content(shape)
                    if shape_data:
                        page_content["raw_content"].append(shape_data)

                        # Extract text content
                        if shape_data["type"] == "text" and shape_data["text"].strip():
                            page_content["text_content"].append(
                                shape_data["text"].strip()
                            )

                        # Categorize content
                        if shape_data["type"] == "image":
                            page_content["images"].append(shape_data)
                        elif shape_data["type"] == "table":
                            page_content["tables"].append(shape_data)

                # Extract slide notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        page_content["notes"] = notes_slide.notes_text_frame.text

                # Only add page if it has content
                if (
                    page_content["title"]
                    or page_content["text_content"]
                    or page_content["images"]
                    or page_content["tables"]
                ):
                    pages_content[page_key] = page_content

            # Calculate total word count
            total_words = 0
            for page in pages_content.values():
                for text in page["text_content"]:
                    total_words += len(text.split())
                if page["notes"]:
                    total_words += len(page["notes"].split())

            # Compile final result
            result = {
                "type": "pptx",
                "metadata": metadata,
                "content": pages_content,
                "summary": {
                    "total_pages": len(pages_content),
                    "pages_with_images": len(
                        [p for p in pages_content.values() if p["images"]]
                    ),
                    "pages_with_tables": len(
                        [p for p in pages_content.values() if p["tables"]]
                    ),
                    "pages_with_notes": len(
                        [p for p in pages_content.values() if p["notes"]]
                    ),
                    "total_text_blocks": sum(
                        len(p["text_content"]) for p in pages_content.values()
                    ),
                },
                "total_word_count": total_words,
            }

            return result

        except Exception as e:
            logger.error(f"Error parsing PPTX {filename}: {str(e)}")
            raise Exception(f"Failed to parse PPTX: {str(e)}")

    @staticmethod
    async def _parse_docx(file_content: bytes, filename: str) -> Dict[str, Any]:
        """Parse Word document content with heading-based organization."""
        try:
            docx_file = io.BytesIO(file_content)
            document = DocxDocument(docx_file)

            # Extract document metadata
            metadata = {
                "filename": filename,
                "file_type": "DOCX",
                "parsed_at": datetime.utcnow().isoformat(),
                "total_paragraphs": len(document.paragraphs),
                "total_tables": len(document.tables),
                "total_images": len(document.inline_shapes),
                "title": document.core_properties.title or "",
                "author": document.core_properties.author or "",
                "subject": document.core_properties.subject or "",
                "created": (
                    document.core_properties.created.isoformat()
                    if document.core_properties.created
                    else ""
                ),
                "modified": (
                    document.core_properties.modified.isoformat()
                    if document.core_properties.modified
                    else ""
                ),
            }

            # Organize content by headings
            content_by_headings = {}
            current_heading = None
            current_subheading = None
            content_buffer = []

            for paragraph in document.paragraphs:
                if not paragraph.text.strip():
                    continue

                # Check if it's a heading
                style_name = paragraph.style.name if paragraph.style else "Normal"
                is_heading = style_name.startswith("Heading")

                if is_heading:
                    # Save previous content if exists
                    if current_heading and content_buffer:
                        if current_subheading:
                            if current_heading not in content_by_headings:
                                content_by_headings[current_heading] = {}
                            content_by_headings[current_heading][
                                current_subheading
                            ] = content_buffer.copy()
                        else:
                            content_by_headings[current_heading] = content_buffer.copy()
                        content_buffer = []

                    # Determine heading level
                    if style_name == "Heading 1":
                        current_heading = paragraph.text.strip()
                        current_subheading = None
                    elif style_name == "Heading 2":
                        current_subheading = paragraph.text.strip()
                    else:
                        # Handle other heading levels as subheadings
                        if current_subheading:
                            current_subheading = (
                                f"{current_subheading} - {paragraph.text.strip()}"
                            )
                        else:
                            current_subheading = paragraph.text.strip()
                else:
                    # Regular paragraph - add to content buffer
                    para_content = {
                        "text": paragraph.text.strip(),
                        "style": style_name,
                        "formatting": DocumentParser._extract_paragraph_formatting(
                            paragraph
                        ),
                    }
                    content_buffer.append(para_content)

            # Save final content
            if current_heading and content_buffer:
                if current_subheading:
                    if current_heading not in content_by_headings:
                        content_by_headings[current_heading] = {}
                    content_by_headings[current_heading][
                        current_subheading
                    ] = content_buffer.copy()
                else:
                    content_by_headings[current_heading] = content_buffer.copy()

            # Extract tables
            tables_data = []
            for table_num, table in enumerate(document.tables, 1):
                table_data = {
                    "type": "table",
                    "table_number": table_num,
                    "rows": len(table.rows),
                    "columns": len(table.columns) if table.rows else 0,
                    "data": [],
                }

                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data["data"].append(row_data)

                tables_data.append(table_data)

            # Extract images
            images_data = []
            for i, shape in enumerate(document.inline_shapes, 1):
                image_data = {
                    "type": "image",
                    "image_number": i,
                    "shape_type": str(shape.type),
                    "width": shape.width.inches if shape.width else None,
                    "height": shape.height.inches if shape.height else None,
                }
                images_data.append(image_data)

            # Calculate total word count
            total_words = 0
            for heading, sections in content_by_headings.items():
                if isinstance(sections, dict):
                    for subheading, content in sections.items():
                        for item in content:
                            if isinstance(item, dict) and "text" in item:
                                total_words += len(item["text"].split())
                elif isinstance(sections, list):
                    for item in sections:
                        if isinstance(item, dict) and "text" in item:
                            total_words += len(item["text"].split())

            for table in tables_data:
                for row in table["data"]:
                    for cell in row:
                        total_words += len(cell.split())

            # Compile final result
            result = {
                "type": "docx",
                "metadata": metadata,
                "content": content_by_headings,
                "tables": tables_data,
                "images": images_data,
                "summary": {
                    "total_headings": len(content_by_headings),
                    "total_sections": sum(
                        len(v) if isinstance(v, dict) else 1
                        for v in content_by_headings.values()
                    ),
                    "total_paragraphs": sum(
                        (
                            len(v)
                            if isinstance(v, list)
                            else sum(
                                len(sub) if isinstance(sub, list) else 0
                                for sub in v.values()
                            )
                        )
                        for v in content_by_headings.values()
                    ),
                    "total_tables": len(tables_data),
                    "total_images": len(images_data),
                },
                "total_word_count": total_words,
            }

            return result

        except Exception as e:
            logger.error(f"Error parsing DOCX {filename}: {str(e)}")
            raise Exception(f"Failed to parse DOCX: {str(e)}")

    @staticmethod
    async def parse_website(url: str) -> WebsiteContent:
        """
        Parse website content and extract text.

        Args:
            url: Website URL to parse

        Returns:
            WebsiteContent object with extracted data
        """
        try:
            # Validate URL
            if not url.startswith(("http://", "https://")):
                url = "https://" + url

            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=30) as response:
                    if response.status != 200:
                        raise Exception(
                            f"Failed to fetch website: HTTP {response.status}"
                        )

                    html_content = await response.text()

            # Parse HTML content
            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Extract title
            title = soup.find("title")
            title_text = title.get_text().strip() if title else ""

            # Extract main content
            # Try to find main content areas
            main_content = ""

            # Look for common content containers
            content_selectors = [
                "main",
                "article",
                ".content",
                "#content",
                ".main-content",
                "#main-content",
                ".post-content",
                ".entry-content",
                ".article-content",
            ]

            for selector in content_selectors:
                content_elem = soup.select_one(selector)
                if content_elem:
                    main_content = content_elem.get_text(separator=" ", strip=True)
                    break

            # If no main content found, get all text
            if not main_content:
                main_content = soup.get_text(separator=" ", strip=True)

            # Clean up text
            main_content = " ".join(main_content.split())

            # Extract metadata
            metadata = {
                "url": url,
                "title": title_text,
                "description": "",
                "keywords": "",
                "author": "",
                "word_count": len(main_content.split()),
                "extracted_at": datetime.utcnow().isoformat(),
            }

            # Try to extract meta description
            meta_desc = soup.find("meta", attrs={"name": "description"})
            if meta_desc:
                metadata["description"] = meta_desc.get("content", "")

            # Try to extract meta keywords
            meta_keywords = soup.find("meta", attrs={"name": "keywords"})
            if meta_keywords:
                metadata["keywords"] = meta_keywords.get("content", "")

            # Try to extract author
            meta_author = soup.find("meta", attrs={"name": "author"})
            if meta_author:
                metadata["author"] = meta_author.get("content", "")

            return WebsiteContent(
                url=url, title=title_text, content=main_content, metadata=metadata
            )

        except Exception as e:
            logger.error(f"Error parsing website {url}: {str(e)}")
            raise Exception(f"Failed to parse website: {str(e)}")

    @staticmethod
    def _extract_shape_content(shape) -> Optional[Dict[str, Any]]:
        """Extract content from a shape in PPTX"""
        try:
            shape_data = {"type": "unknown", "text": "", "position": {}, "size": {}}

            # Get position and size
            if hasattr(shape, "left") and hasattr(shape, "top"):
                shape_data["position"] = {"left": shape.left, "top": shape.top}

            if hasattr(shape, "width") and hasattr(shape, "height"):
                shape_data["size"] = {"width": shape.width, "height": shape.height}

            # Handle different shape types
            if shape.has_text_frame:
                shape_data["type"] = "text"
                shape_data["text"] = shape.text

                # Extract text formatting
                if shape.text_frame:
                    shape_data["paragraphs"] = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            shape_data["paragraphs"].append(
                                {
                                    "text": para_text,
                                    "alignment": (
                                        str(paragraph.alignment)
                                        if paragraph.alignment
                                        else None
                                    ),
                                    "level": paragraph.level,
                                }
                            )

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_data["type"] = "image"
                shape_data["text"] = (
                    f"[Image: {shape.name if hasattr(shape, 'name') else 'Unknown'}]"
                )

            elif shape.has_table:
                shape_data["type"] = "table"
                table = shape.table
                shape_data["table_data"] = {
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": [],
                }

                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    shape_data["table_data"]["data"].append(row_data)

            else:
                shape_data["type"] = "shape"
                shape_data["text"] = shape.name if hasattr(shape, "name") else ""

            return (
                shape_data
                if shape_data["text"] or shape_data["type"] in ["image", "table"]
                else None
            )

        except Exception as e:
            logger.warning(f"Error extracting shape content: {str(e)}")
            return None

    @staticmethod
    def _extract_paragraph_formatting(paragraph) -> Dict[str, Any]:
        """Extract formatting information from a paragraph"""
        formatting = {"runs": []}

        for run in paragraph.runs:
            if run.text.strip():
                run_data = {
                    "text": run.text,
                    "bold": run.bold,
                    "italic": run.italic,
                    "underline": run.underline,
                    "font_size": str(run.font.size) if run.font.size else None,
                    "font_name": run.font.name if run.font.name else None,
                }
                formatting["runs"].append(run_data)

        return formatting

    @staticmethod
    def save_content_to_json(content: Dict[str, Any], file_path: str) -> str:
        """
        Save parsed content to JSON file.

        Args:
            content: Parsed content dictionary
            file_path: Path where to save the JSON file

        Returns:
            Path to the saved JSON file
        """
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(content, f, indent=2, ensure_ascii=False)
            return file_path
        except Exception as e:
            logger.error(f"Error saving content to JSON {file_path}: {str(e)}")
            raise Exception(f"Failed to save content to JSON: {str(e)}")
