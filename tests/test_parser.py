"""
Document parsing service: PDF, PPTX, DOCX.
- PDF: page-level extraction, heading detection, text cleaning, LLM-friendly fields
- PPTX/DOCX: structured extraction consistent with prior behavior
"""

import io
import json
import math
import os
import re
import tempfile
import unicodedata
from collections import defaultdict
from contextlib import ExitStack
from dataclasses import dataclass
from datetime import datetime
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:  # Optional dependencies, best-effort imports
    import fitz  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    fitz = None

try:
    import pdfplumber  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    pdfplumber = None

try:
    import camelot  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    camelot = None

try:
    import tabula  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    tabula = None

try:
    import pytesseract  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None

try:
    from PIL import Image  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    Image = None

from logger import logger
from models.document import DocumentType

# ========================= Helpers: headings (PDF) =========================
HEADING_LEXICON = {
    "chapter": [
        "chapter",
        "unit",
        "module",
        "part",
        "lesson",
        "week",
        "lecture",
    ],
    "topic": [
        "section",
        "topic",
        "theme",
        "lesson",
    ],
    "answer_key": [
        "answer key",
        "answers",
        "solutions",
        "selected answers",
        "odd answers",
        "even answers",
    ],
}

CANONICAL_CHAPTER_LABELS = {"chapter", "unit", "module", "part"}
SYNONYM_CHAPTER_LABELS = set(label for label in HEADING_LEXICON["chapter"] if label not in CANONICAL_CHAPTER_LABELS)
ALL_CHAPTER_LABELS = CANONICAL_CHAPTER_LABELS | SYNONYM_CHAPTER_LABELS
CHAPTER_KEYWORDS = tuple(sorted(ALL_CHAPTER_LABELS))
TOPIC_KEYWORDS = tuple(sorted({*HEADING_LEXICON["topic"], "module", "unit"}))
CHAPTER_LINE_RE = re.compile(
    r"^\s*(?:" + "|".join(CHAPTER_KEYWORDS) + r")\s*(\d{1,3})\b[\s:.\-|\u2013\u2014\)]*\s*(.*)$",
    re.IGNORECASE,
)
CHAP_SYNONYM_RE = re.compile(
    r"^\s*(?P<label>" + "|".join(sorted(ALL_CHAPTER_LABELS)) + r")\s+(?P<num>\d{1,3})\b[\s:.\-|\u2013\u2014\)]*\s*(?P<title>.*)$",
    re.IGNORECASE,
)
TOPIC_LINE_RE = re.compile(
    r"^\s*(?:" + "|".join(TOPIC_KEYWORDS) + r")?\s*(\d{1,3}(?:\s*[.\-]\s*\d+)+)\b[\s:.\-|\u2013\u2014\)]*\s*(.*)$",
    re.IGNORECASE,
)
TOPIC_NUM_RE = re.compile(
    r"^\s*(?:section|topic|lesson|theme)?\s*(\d{1,3}(?:\s*\.\s*\d+)+)\b[:.\-\)]?\s*(.*)$",
    re.IGNORECASE,
)
TRAILING_PAGE_RE = re.compile(r"\s+(pp?\.\s*\d+(-\d+)?|\d{1,4})\s*$", re.IGNORECASE)
TOPIC_NUM_NORMALIZER = re.compile(r"(?<!\d)(\d{1,2})\s*[.\-•·]\s*(\d{1,2})(\s*[.\-•·]\s*(\d{1,2}))?(?!\d)")
HEADING_BLACKLIST_PREFIXES = (
    "figure",
    "table",
    "map",
    "click and explore",
    "defining american",
    "key terms",
    "summary",
    "learning objectives",
    "my story",
)

FIGURE_CAPTION_RE = re.compile(
    r"^(?:fig(?:ure)?)\s*[\.\s]*(\d+(?:\.\d+)*)\s*[:\-\. ]\s*(.+)$",
    re.IGNORECASE,
)
TABLE_CAPTION_RE = re.compile(
    r"^(?:table)\s*[\.\s]*(\d+(?:\.\d+)*)\s*[:\-\. ]\s*(.+)$",
    re.IGNORECASE,
)

LEARNING_OBJECTIVE_HEADERS = ("learning objectives", "objectives", "learning goals")
KEY_TERM_HEADERS = ("key terms", "key vocabulary", "glossary")
SIDEBAR_HEADERS = ("case study", "key concept", "focus", "spotlight")

FIGURE_MODE_CHOICES = {"off", "caption", "image", "both"}
TABLE_MODE_CHOICES = {"off", "lines", "camelot", "tabula"}
OCR_THRESHOLDS = {"LOW": 25, "MED": 60, "HIGH": 120}
HEADING_LEXICON_CHOICES = {"auto", "lecture"}

INLINE_TOPIC_KEYWORDS = (
    "how to",
    "key terms",
    "learning objectives",
    "overview",
    "summary",
    "introduction",
    "suggested reading",
    "further reading",
    "case study",
    "tips",
    "checklist",
)

ASSESSMENT_SECTION_HEADERS = (
    "exercises",
    "review questions",
    "practice",
    "assessment",
    "quiz",
    "chapter review",
    "problems",
    "self-check",
    "self check",
)

ANSWER_SECTION_HEADERS = tuple(HEADING_LEXICON["answer_key"])
QUESTION_START_RE = re.compile(r"^\s*(\d{1,3})[\.\)]\s+(.*)")
CHOICE_LINE_RE = re.compile(r"^\s*([A-E])[\.\)]\s+(.*)")
COMPACT_KEY_RE = re.compile(r"(\d{1,3})\s*[-:\.)]?\s*([A-E]|True|False|T|F)", re.IGNORECASE)
TRUE_FALSE_TOKEN_RE = re.compile(r"^\s*(true|false|t|f)\b", re.IGNORECASE)


@dataclass(frozen=True)
class PDFParseOptions:
    figure_mode: str = "both"
    tables_mode: str = "off"
    ocr_threshold: str = "MED"
    heading_lexicon: str = "auto"
    heading_synonyms: bool = True
    toc_window: int = 60
    assessments: bool = True
    answer_keys: bool = True

    @classmethod
    def from_dict(cls, data: Optional[Dict[str, Any]]) -> "PDFParseOptions":
        if not data:
            return cls()
        figure_mode = str(data.get("figure_mode", cls.figure_mode)).lower()
        if figure_mode not in FIGURE_MODE_CHOICES:
            figure_mode = cls.figure_mode
        tables_mode = str(data.get("tables_mode", cls.tables_mode)).lower()
        if tables_mode not in TABLE_MODE_CHOICES:
            tables_mode = cls.tables_mode
        ocr_threshold = str(data.get("ocr_threshold", cls.ocr_threshold)).upper()
        if ocr_threshold not in OCR_THRESHOLDS:
            ocr_threshold = cls.ocr_threshold
        heading_lexicon = str(data.get("heading_lexicon", cls.heading_lexicon)).lower()
        if heading_lexicon not in HEADING_LEXICON_CHOICES:
            heading_lexicon = cls.heading_lexicon
        heading_synonyms_val = data.get("heading_synonyms", cls.heading_synonyms)
        if isinstance(heading_synonyms_val, str):
            heading_synonyms = heading_synonyms_val.lower() == "on"
        else:
            heading_synonyms = bool(heading_synonyms_val)
        assessments_val = data.get("assessments", cls.assessments)
        if isinstance(assessments_val, str):
            assessments = assessments_val.lower() == "on"
        else:
            assessments = bool(assessments_val)
        answer_keys_val = data.get("answer_keys", cls.answer_keys)
        if isinstance(answer_keys_val, str):
            answer_keys = answer_keys_val.lower() == "on"
        else:
            answer_keys = bool(answer_keys_val)
        try:
            toc_window = int(data.get("toc_window", cls.toc_window))
        except (TypeError, ValueError):
            toc_window = cls.toc_window
        toc_window = max(10, toc_window)
        return cls(
            figure_mode=figure_mode,
            tables_mode=tables_mode,
            ocr_threshold=ocr_threshold,
            heading_lexicon=heading_lexicon,
            heading_synonyms=heading_synonyms,
            toc_window=toc_window,
            assessments=assessments,
            answer_keys=answer_keys,
        )


def _nfkc(s: str) -> str:
    if not s:
        return ""
    s = unicodedata.normalize("NFKC", s)
    s = s.replace("\u00A0", " ").replace("\u200B", "").replace("\ufeff", "")
    s = re.sub(r"[\u00ad\u2010\u2011\u2012\u2013\u2014]", "-", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _search_simplify(s: str) -> str:
    s = _nfkc(s).lower()
    return re.sub(r"[^a-z0-9]+", " ", s).strip()


def _collapse_intra_word_spaces(s: str) -> str:
    if not s:
        return s

    pattern = re.compile(r"([a-z])\s{1,2}([a-z])")

    def _repl(match: "re.Match[str]") -> str:
        left, right = match.groups()
        return f"{left}{right}"

    return pattern.sub(_repl, s)


def _normalize_heading_text(s: str) -> str:
    if not s:
        return ""
    text = _nfkc(s)
    text = _collapse_intra_word_spaces(text)
    text = re.sub(r"(?:\b[A-Z]\s+){2,}\b[A-Z]\b", lambda m: m.group(0).replace(" ", ""), text)
    text = re.sub(r"(\d)\s*[.\-•·]\s*(\d+)", r"\1.\2", text)
    text = re.sub(r"(\d)\s*\.\s*(\d)\s*\.\s*(\d+)", r"\1.\2.\3", text)
    text = TRAILING_PAGE_RE.sub("", text).strip()
    return text


def _normalize_topic_number(raw: str) -> Optional[str]:
    if not raw:
        return None
    match = TOPIC_NUM_NORMALIZER.search(raw)
    if not match:
        return None
    parts = [match.group(1), match.group(2)]
    if match.group(4):
        parts.append(match.group(4).strip())
    normalized = ".".join(part.lstrip("0") or "0" for part in parts)
    return normalized


def _extract_first_int(value: Any) -> Optional[int]:
    if value is None:
        return None
    if isinstance(value, int):
        return value
    try:
        text = str(value)
    except Exception:
        return None
    match = re.search(r"\d+", text)
    if not match:
        return None
    try:
        return int(match.group(0))
    except ValueError:
        return None


def _compute_char_ngrams(text: str, n: int = 3) -> List[str]:
    cleaned = re.sub(r"\s+", " ", text or "").lower().strip()
    if len(cleaned) < n:
        return [cleaned] if cleaned else []
    return [cleaned[i : i + n] for i in range(len(cleaned) - n + 1)]


def _score_topic_candidate(page_text: str, candidate_title: str) -> float:
    if not page_text or not candidate_title:
        return 0.0
    ngrams_page = set(_compute_char_ngrams(page_text[:500]))
    ngrams_title = set(_compute_char_ngrams(candidate_title))
    if not ngrams_page or not ngrams_title:
        overlap = 0.0
    else:
        overlap = len(ngrams_page & ngrams_title) / len(ngrams_page | ngrams_title)

    tokens_page = re.findall(r"[A-Za-z0-9]+", page_text[:1000].lower())
    tokens_title = set(re.findall(r"[A-Za-z0-9]+", candidate_title.lower()))
    keyword_hits = sum(1 for token in tokens_page if token in tokens_title)
    keyword_score = min(1.0, keyword_hits / 5.0)
    return round(0.7 * overlap + 0.3 * keyword_score, 4)


BBox = Tuple[float, float, float, float]


def _coerce_bbox(raw: Sequence[float]) -> BBox:
    if len(raw) != 4:
        return (0.0, 0.0, 0.0, 0.0)
    return tuple(float(v) for v in raw)  # type: ignore[return-value]


def _round_bbox(bbox: BBox, precision: int = 2) -> List[float]:
    return [round(coord, precision) for coord in bbox]


def _same_column(b1: BBox, b2: BBox, tolerance: float = 35.0) -> bool:
    overlap = min(b1[2], b2[2]) - max(b1[0], b2[0])
    if overlap > 0:
        return True
    return abs(b1[0] - b2[0]) <= tolerance or abs(b1[2] - b2[2]) <= tolerance


def _vertical_gap(top_box: BBox, bottom_box: BBox) -> float:
    if top_box[1] >= bottom_box[3]:
        return top_box[1] - bottom_box[3]
    if bottom_box[1] >= top_box[3]:
        return bottom_box[1] - top_box[3]
    return 0.0


def _bbox_area(bbox: BBox) -> float:
    return max(0.0, (bbox[2] - bbox[0]) * (bbox[3] - bbox[1]))


def _split_toc_line(line: str) -> Tuple[str, Optional[int]]:
    cleaned = _nfkc(line).strip()
    if not cleaned:
        return "", None
    page_candidate = re.search(r"(\d{1,4})\s*$", cleaned)
    if page_candidate:
        try:
            page_number = int(page_candidate.group(1))
        except ValueError:
            page_number = None
        title = cleaned[: page_candidate.start()].strip()
        return title, page_number
    return cleaned, None


def _format_synonym_label(label: str, number: str) -> str:
    return f"{label.strip().title()} {number.strip()}"


def _collect_subsection_lines(lines: Sequence[str], start_index: int) -> List[str]:
    collected: List[str] = []
    for raw in lines[start_index + 1 :]:
        current = _nfkc(raw).strip()
        if not current:
            if collected:
                break
            continue
        if CHAPTER_LINE_RE.match(current) or TOPIC_LINE_RE.match(current):
            break
        collected.append(current)
        if len(collected) >= 25:
            break
    return collected


def _extract_learning_objectives(lines: Sequence[str]) -> List[str]:
    objectives: List[str] = []
    collecting = False
    for line in lines:
        norm = _nfkc(line).strip()
        lower = norm.lower()
        if any(lower.startswith(header) for header in LEARNING_OBJECTIVE_HEADERS):
            collecting = True
            continue
        if collecting:
            if not norm:
                break
            if norm.isupper() and len(norm.split()) <= 6:
                break
            objectives.append(norm)
    return objectives


def _extract_key_terms(lines: Sequence[str]) -> List[str]:
    terms: List[str] = []
    collecting = False
    for line in lines:
        norm = _nfkc(line).strip()
        lower = norm.lower()
        if any(lower.startswith(header) for header in KEY_TERM_HEADERS):
            collecting = True
            continue
        if collecting:
            if not norm:
                collecting = False
                continue
            if norm.isupper() and len(norm.split()) <= 6:
                collecting = False
                continue
            fragments = re.split(r"[,;]\s*", norm)
            for frag in fragments:
                cleaned = frag.strip()
                if cleaned and cleaned not in terms:
                    terms.append(cleaned)
    return terms


def _extract_footnotes(lines: Sequence[str]) -> Tuple[List[str], List[str]]:
    footnotes: List[str] = []
    references: List[str] = []
    tail = list(lines[-8:]) if lines else []
    for line in tail:
        norm = _nfkc(line).strip()
        if not norm:
            continue
        if re.match(r"^\(?\d+[)\.\]]", norm):
            footnotes.append(norm)
        elif re.search(r"\(\d{4}[a-z]?\)", norm) or re.search(r"https?://", norm.lower()):
            references.append(norm)
    return footnotes, references


def _extract_sidebars(lines: Sequence[str]) -> List[str]:
    sidebars: List[str] = []
    for line in lines:
        norm = _nfkc(line).strip()
        lower = norm.lower()
        if any(lower.startswith(header) for header in SIDEBAR_HEADERS):
            sidebars.append(norm)
    return sidebars


def _normalize_answer_token(token: str) -> str:
    cleaned = _nfkc(token).strip()
    if not cleaned:
        return ""
    upper = cleaned.upper()
    if upper in {"TRUE", "T"}:
        return "True"
    if upper in {"FALSE", "F"}:
        return "False"
    if len(upper) == 1 and upper.isalpha():
        return upper
    return cleaned


def _detect_answer_key_on_page(lines: Sequence[str]) -> Optional[Dict[str, Any]]:
    if not lines:
        return None

    header_index: Optional[int] = None
    header_text: Optional[str] = None
    odd_only = False

    for idx, raw in enumerate(lines[:20]):
        normalized = _normalize_heading_text(raw)
        lower = normalized.lower()
        if any(lower.startswith(header) for header in ANSWER_SECTION_HEADERS) or "answer key" in lower or "solutions" in lower:
            header_index = idx
            header_text = normalized or raw.strip()
            odd_only = "odd" in lower and "even" not in lower
            break

    if header_index is None:
        return None

    entries: List[Dict[str, Any]] = []
    answer_lines: List[str] = []
    seen_lines: set[str] = set()

    for raw in lines[header_index:]:
        norm_line = _nfkc(raw).strip()
        if not norm_line:
            continue
        lowered = norm_line.lower()
        if any(lowered.startswith(header) for header in ASSESSMENT_SECTION_HEADERS):
            break

        line_numbers: set[str] = set()
        match = QUESTION_START_RE.match(norm_line)
        if match:
            number = match.group(1)
            remainder = match.group(2).strip()
            answer_token: Optional[str] = None
            tf_match = TRUE_FALSE_TOKEN_RE.match(remainder)
            if tf_match:
                answer_token = tf_match.group(1)
            else:
                letter_match = re.search(r"\b([A-E])\b", remainder, re.IGNORECASE)
                if letter_match:
                    answer_token = letter_match.group(1)
            if answer_token:
                entries.append({
                    "number": number,
                    "answer": _normalize_answer_token(answer_token),
                    "raw": norm_line,
                })
                line_numbers.add(number)

        compact_hits = COMPACT_KEY_RE.findall(norm_line)
        for num, ans in compact_hits:
            if num in line_numbers:
                continue
            entries.append({
                "number": num,
                "answer": _normalize_answer_token(ans),
                "raw": norm_line,
            })
            line_numbers.add(num)

        if line_numbers and norm_line not in seen_lines:
            answer_lines.append(norm_line)
            seen_lines.add(norm_line)

    if not entries:
        return None

    return {
        "title": header_text or "Answer Key",
        "entries": entries,
        "lines": answer_lines,
        "odd_only": odd_only,
        "header_index": header_index,
    }


def _extract_assessment_blocks(lines: Sequence[str]) -> List[Dict[str, Any]]:
    if not lines:
        return []

    blocks: List[Dict[str, Any]] = []
    current_block: Optional[Dict[str, Any]] = None
    current_question: Optional[Dict[str, Any]] = None

    def _finalize_question(question: Dict[str, Any]) -> Dict[str, Any]:
        prompt_lines = [segment for segment in question.get("prompt_lines", []) if segment]
        prompt = " ".join(prompt_lines).strip()
        return {
            "number": question.get("number"),
            "prompt": prompt,
            "choices": question.get("choices", []),
        }

    def _finalize_block(block: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        questions = [q for q in block.get("questions", []) if q.get("number")]
        if not questions:
            return None
        block_copy = {
            "title": block.get("title"),
            "questions": questions,
        }
        if block.get("instructions"):
            block_copy["instructions"] = block["instructions"]
        return block_copy

    for raw in lines:
        norm = _nfkc(raw).strip()
        if not norm:
            continue
        lower = norm.lower()

        if any(lower.startswith(header) for header in ANSWER_SECTION_HEADERS):
            if current_question and current_block:
                finalized = _finalize_question(current_question)
                if finalized:
                    current_block.setdefault("questions", []).append(finalized)
            if current_block:
                finalized_block = _finalize_block(current_block)
                if finalized_block:
                    blocks.append(finalized_block)
            current_block = None
            current_question = None
            continue

        if any(lower.startswith(header) for header in ASSESSMENT_SECTION_HEADERS):
            if current_question and current_block:
                finalized = _finalize_question(current_question)
                if finalized:
                    current_block.setdefault("questions", []).append(finalized)
            if current_block:
                finalized_block = _finalize_block(current_block)
                if finalized_block:
                    blocks.append(finalized_block)
            current_block = {"title": _normalize_heading_text(norm), "questions": []}
            current_question = None
            continue

        if current_block is None:
            continue

        q_match = QUESTION_START_RE.match(norm)
        if q_match:
            if current_question:
                finalized = _finalize_question(current_question)
                if finalized:
                    current_block.setdefault("questions", []).append(finalized)
            current_question = {
                "number": q_match.group(1),
                "prompt_lines": [q_match.group(2).strip()] if q_match.group(2).strip() else [],
                "choices": [],
            }
            continue

        if current_question is None:
            current_block.setdefault("instructions", []).append(norm)
            continue

        choice_match = CHOICE_LINE_RE.match(norm)
        if choice_match:
            letter = choice_match.group(1)
            text = choice_match.group(2).strip()
            formatted = f"{letter}. {text}" if text else letter
            current_question.setdefault("choices", []).append(formatted)
            continue

        if TRUE_FALSE_TOKEN_RE.match(norm) and not current_question.get("choices"):
            current_question.setdefault("choices", []).append(norm)
            continue

        current_question.setdefault("prompt_lines", []).append(norm)

    if current_question and current_block:
        finalized = _finalize_question(current_question)
        if finalized:
            current_block.setdefault("questions", []).append(finalized)
    if current_block:
        finalized_block = _finalize_block(current_block)
        if finalized_block:
            blocks.append(finalized_block)

    return blocks


def _build_question_ref(page: Dict[str, Any], question_number: str) -> str:
    topic = page.get("topic") or ""
    chapter = page.get("chapter") or ""
    topic_clean = str(topic).strip()
    if topic_clean:
        return f"{topic_clean}-Q{question_number}"
    chapter_num = _extract_first_int(chapter)
    if chapter_num is not None:
        return f"{chapter_num}-Q{question_number}"
    return f"Q{question_number}"


def _aggregate_assessment_data(
    pages: Dict[str, Dict[str, Any]],
    options: PDFParseOptions,
    diagnostics: Dict[str, Any],
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    diagnostics.setdefault("answer_key_pages_found", [])
    diagnostics.setdefault("assessment_blocks_found", [])
    diagnostics.setdefault("orphan_answers", [])
    diagnostics.setdefault("unmatched_questions", [])

    assessments: List[Dict[str, Any]] = []
    answer_keys: List[Dict[str, Any]] = []

    if not options.assessments and not options.answer_keys:
        for page in pages.values():
            page.setdefault("answer_key_lines", [])
            page.setdefault("assessment_blocks", [])
        return assessments, answer_keys

    ordered_pages = sorted(pages.values(), key=lambda p: p.get("page_number") or 0)

    for page in ordered_pages:
        page.setdefault("answer_key_lines", [])
        page.setdefault("assessment_blocks", [])

    answer_map: Dict[str, List[Dict[str, Any]]] = {}
    assessment_index: Dict[str, List[Dict[str, Any]]] = {}

    current_answer_section: Optional[Dict[str, Any]] = None
    previous_answer_page: Optional[int] = None

    for page in ordered_pages:
        pn = page.get("page_number") or 0
        lines = page.get("lines") or []

        answer_detection = _detect_answer_key_on_page(lines) if options.answer_keys else None

        if answer_detection:
            existing_type = page.get("heading_type")
            if existing_type != "answer_key":
                logger.info("auto-repair: page %s heading reclassified to answer_key", pn)
            page["heading_type"] = "answer_key"
            page["answer_key_lines"] = answer_detection.get("lines", [])
            diagnostics["answer_key_pages_found"].append(pn)

            title = answer_detection.get("title") or "Answer Key"
            contiguous = previous_answer_page is not None and pn == previous_answer_page + 1
            same_title = current_answer_section and current_answer_section.get("title") == title
            if not current_answer_section or not contiguous or not same_title:
                if current_answer_section:
                    output_section = {
                        "title": current_answer_section["title"],
                        "start_page": current_answer_section["start_page"],
                        "end_page": current_answer_section["end_page"],
                        "mapping": current_answer_section["mapping"],
                    }
                    answer_keys.append(output_section)
                current_answer_section = {
                    "title": title,
                    "start_page": pn,
                    "end_page": pn,
                    "mapping": [],
                    "odd_only": answer_detection.get("odd_only", False),
                }
            else:
                current_answer_section["end_page"] = pn
                current_answer_section["odd_only"] = current_answer_section.get("odd_only", False) or answer_detection.get("odd_only", False)

            for entry in answer_detection.get("entries", []):
                number = entry.get("number")
                answer_value = entry.get("answer")
                if not number or not answer_value:
                    continue
                ref = _build_question_ref(page, str(number))
                answer_record = {
                    "ref": ref,
                    "answer": answer_value,
                    "page": pn,
                    "odd_only": answer_detection.get("odd_only", False),
                }
                current_answer_section["mapping"].append({"ref": ref, "answer": answer_value})
                answer_map.setdefault(ref, []).append(answer_record)

            previous_answer_page = pn
        else:
            if current_answer_section:
                output_section = {
                    "title": current_answer_section["title"],
                    "start_page": current_answer_section["start_page"],
                    "end_page": current_answer_section["end_page"],
                    "mapping": current_answer_section["mapping"],
                }
                answer_keys.append(output_section)
                current_answer_section = None
                previous_answer_page = None

        blocks = _extract_assessment_blocks(lines) if options.assessments else []
        if blocks:
            page["assessment_blocks"] = blocks
            diagnostics["assessment_blocks_found"].append({
                "page": pn,
                "blocks": len(blocks),
                "questions": sum(len(block.get("questions", [])) for block in blocks),
            })
            for block in blocks:
                for question in block.get("questions", []):
                    number = question.get("number")
                    if not number:
                        continue
                    ref = _build_question_ref(page, str(number))
                    entry = {
                        "id": ref,
                        "prompt": question.get("prompt", ""),
                        "choices": question.get("choices", []),
                        "answer": None,
                        "source_page": pn,
                    }
                    assessments.append(entry)
                    assessment_index.setdefault(ref, []).append(entry)

    if current_answer_section:
        output_section = {
            "title": current_answer_section["title"],
            "start_page": current_answer_section["start_page"],
            "end_page": current_answer_section["end_page"],
            "mapping": current_answer_section["mapping"],
        }
        answer_keys.append(output_section)

    matched_answers: set[str] = set()

    for ref, questions in assessment_index.items():
        answers = answer_map.get(ref)
        if not answers:
            diagnostics["unmatched_questions"].append(ref)
            continue
        answer_choice = answers[0]
        matched_answers.add(ref)
        for question in questions:
            question["answer"] = answer_choice.get("answer")
            if answer_choice.get("odd_only"):
                question["answer_source"] = "odd_only"

    for ref, records in answer_map.items():
        if ref not in matched_answers:
            diagnostics["orphan_answers"].extend(records)

    return assessments, answer_keys


def _categorize_word_count(word_count: int) -> str:
    if word_count >= 200:
        return "high"
    if word_count >= 80:
        return "medium"
    return "low"


def _topic_tuple(topic: str) -> Tuple[int, ...]:
    numeric = []
    for part in re.split(r"[.\-]\s*", topic.strip()):
        try:
            numeric.append(int(part))
        except ValueError:
            continue
    return tuple(numeric)


def _run_ocr_on_page(doc_page: Any) -> str:
    if doc_page is None or pytesseract is None or Image is None or fitz is None:
        return ""
    try:
        matrix = fitz.Matrix(2, 2)  # type: ignore[attr-defined]
        pix = doc_page.get_pixmap(matrix=matrix)
        mode = "RGB" if pix.alpha == 0 else "RGBA"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        return pytesseract.image_to_string(img)
    except Exception as exc:  # pragma: no cover - OCR fallback best effort
        logger.debug("OCR fallback failed on page %s: %s", getattr(doc_page, "number", "?"), exc)
        return ""


def _auto_repair_topic_assignment(
    page: Dict[str, Any],
    pn: int,
    toc: Dict[str, Dict[str, Any]],
    chapter_ranges: Dict[str, Tuple[int, int]],
    topic_ranges: Dict[str, Tuple[int, int]],
    diagnostics: Dict[str, Any],
    reason: str,
) -> Optional[str]:
    diagnostics.setdefault("auto_repairs", [])

    candidate_chapter: Optional[str] = None
    for chapter_num, bounds in chapter_ranges.items():
        start, end = bounds
        if start <= pn <= end:
            candidate_chapter = chapter_num
            break

    if not candidate_chapter or candidate_chapter not in toc:
        return None

    chapter_info = toc[candidate_chapter]
    topics = sorted(chapter_info.get("topics", []), key=lambda t: t.get("start_page") or 0)
    if not topics:
        return None

    candidates = []
    for idx, topic in enumerate(topics):
        start = topic.get("start_page")
        end = topic.get("end_page")
        if start is None or end is None:
            continue
        if start <= pn <= end:
            candidates.append(topic)
            if idx - 1 >= 0:
                candidates.append(topics[idx - 1])
            if idx + 1 < len(topics):
                candidates.append(topics[idx + 1])

    if not candidates:
        # fallback to closest topics by distance
        distances = [
            (abs((topic.get("start_page") or pn) - pn), topic)
            for topic in topics
            if topic.get("start_page") is not None
        ]
        distances.sort(key=lambda item: item[0])
        candidates = [topic for _, topic in distances[:2]]

    best_topic: Optional[str] = None
    best_score = 0.0
    best_title: Optional[str] = None
    page_text = page.get("content_cleaned") or page.get("text") or ""

    seen: set[str] = set()
    for topic in candidates:
        topic_num = topic.get("num")
        if not topic_num or topic_num in seen:
            continue
        seen.add(topic_num)
        title = topic.get("title") or ""
        score = _score_topic_candidate(page_text, title)
        if score >= 0.35 and score > best_score:
            bounds = topic_ranges.get(topic_num)
            if bounds and bounds[0] <= pn <= bounds[1]:
                best_topic = topic_num
                best_score = score
                best_title = title

    if best_topic:
        page["topic"] = best_topic
        if best_title:
            page["topic_title"] = best_title
        diagnostics["auto_repairs"].append({
            "page": pn,
            "from": reason,
            "to": best_topic,
            "score": best_score,
        })
        logger.info("auto-repair: page %s topic reassigned to %s (score=%.2f)", pn, best_topic, best_score)
        return best_topic

    return None


def _detect_visual_elements(
    page_index: int,
    doc_page: Optional[Any],
    layout: Optional[Dict[str, Any]],
    options: PDFParseOptions,
    tables_source: Optional[str] = None,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], List[Dict[str, Any]]]:
    if doc_page is None or layout is None:
        return [], [], []

    figure_entries: List[Dict[str, Any]] = []
    table_entries: List[Dict[str, Any]] = []
    orphan_captions: List[Dict[str, Any]] = []

    images: List[Dict[str, Any]] = []
    text_lines: List[Dict[str, Any]] = []

    for block in layout.get("blocks", []):
        block_type = block.get("type")
        bbox = _coerce_bbox(block.get("bbox", (0.0, 0.0, 0.0, 0.0)))
        if block_type == 1:
            images.append({"bbox": bbox})
        elif block_type == 0:
            for line in block.get("lines", []):
                text = "".join(span.get("text", "") for span in line.get("spans", []))
                cleaned = _nfkc(text).strip()
                if cleaned:
                    text_lines.append({"text": cleaned, "bbox": _coerce_bbox(line.get("bbox", bbox))})

    figure_captions: List[Dict[str, Any]] = []
    table_captions: List[Dict[str, Any]] = []
    for entry in text_lines:
        lower = entry["text"].lower()
        fig_match = FIGURE_CAPTION_RE.match(lower)
        if fig_match:
            figure_captions.append(
                {
                    "number": fig_match.group(1),
                    "title": fig_match.group(2).strip(),
                    "bbox": entry["bbox"],
                    "raw": entry["text"],
                }
            )
            continue
        table_match = TABLE_CAPTION_RE.match(lower)
        if table_match:
            table_captions.append(
                {
                    "number": table_match.group(1),
                    "title": table_match.group(2).strip(),
                    "bbox": entry["bbox"],
                    "raw": entry["text"],
                }
            )

    image_records: List[Dict[str, Any]] = []
    if options.figure_mode in {"image", "both"} and images:
        for idx, image in enumerate(images, 1):
            bbox = image["bbox"]
            entry = {
                "id": f"Fig.{page_index}.{idx}",
                "caption": None,
                "bbox": _round_bbox(bbox),
                "confidence": 0.5,
            }
            figure_entries.append(entry)
            image_records.append({"bbox": bbox, "entry": entry, "assigned": False})

    def _assign_caption(
        captions: Iterable[Dict[str, Any]],
        target_records: List[Dict[str, Any]],
        prefix: str,
        allow_orphans: bool,
    ) -> None:
        for caption in captions:
            matched = False
            best_record: Optional[Dict[str, Any]] = None
            best_distance = math.inf
            if target_records:
                for record in target_records:
                    bbox = record["bbox"]
                    distance = _vertical_gap(caption["bbox"], bbox)
                    if not _same_column(caption["bbox"], bbox):
                        distance += 50.0
                    if distance < best_distance:
                        best_distance = distance
                        best_record = record
                if best_record is not None and best_distance <= 250.0:
                    entry = best_record["entry"]
                    entry["id"] = f"{prefix} {caption['number']}"
                    entry["caption"] = caption["title"] or caption["raw"]
                    confidence = max(0.35, 1.0 - min(best_distance / 300.0, 1.0))
                    entry["confidence"] = round(confidence, 2)
                    best_record["assigned"] = True
                    matched = True
            if not matched:
                if allow_orphans:
                    orphan_captions.append(
                        {
                            "page": page_index,
                            "id": f"{prefix} {caption['number']}",
                            "caption": caption["raw"],
                        }
                    )
                else:
                    figure_entries.append(
                        {
                            "id": f"{prefix} {caption['number']}",
                            "caption": caption["title"] or caption["raw"],
                            "bbox": _round_bbox(caption["bbox"]),
                            "confidence": 0.6,
                        }
                    )

    if options.figure_mode in {"caption", "both"}:
        allow_orphans = options.figure_mode == "both"
        _assign_caption(figure_captions, image_records, "Fig", allow_orphans)

    # Append unmatched images with default ids so they still appear with caption None
    if image_records and options.figure_mode in {"image", "both"}:
        for record in image_records:
            if not record.get("assigned"):
                record["entry"]["confidence"] = round(min(record["entry"].get("confidence", 0.5) + 0.05, 0.75), 2)

    # Table handling
    tables_mode = options.tables_mode
    if tables_mode == "camelot" and camelot is None:
        logger.info("Camelot not available; falling back to caption heuristic for tables")
        tables_mode = "lines"
    elif tables_mode == "tabula" and tabula is None:
        logger.info("tabula-py not available; falling back to caption heuristic for tables")
        tables_mode = "lines"

    if tables_mode == "camelot" and tables_source:
        try:
            camelot_tables = camelot.read_pdf(tables_source, pages=str(page_index))  # type: ignore[attr-defined]
            for idx, cam_table in enumerate(camelot_tables, 1):
                bbox = _coerce_bbox(getattr(cam_table, "_bbox", (0.0, 0.0, 0.0, 0.0)))
                caption_bits = []
                try:
                    header = cam_table.df.iloc[0]
                    caption_bits = [str(val).strip() for val in header if str(val).strip()]
                except Exception:
                    caption_bits = []
                caption_text = "; ".join(caption_bits)
                confidence = cam_table.parsing_report.get("accuracy") if cam_table.parsing_report else None
                confidence_score = 0.7
                if confidence is not None:
                    try:
                        confidence_score = max(0.3, min(float(confidence) / 100.0, 1.0))
                    except Exception:
                        confidence_score = 0.7
                table_entries.append(
                    {
                        "id": f"Table {page_index}.{idx}",
                        "caption": caption_text or None,
                        "bbox": _round_bbox(bbox),
                        "confidence": round(confidence_score, 2),
                    }
                )
        except Exception as camelot_exc:  # pragma: no cover - optional dependency
            logger.debug("Camelot failed on page %s: %s", page_index, camelot_exc)
            tables_mode = "lines"

    if tables_mode == "tabula" and tables_source:
        try:
            tabula_tables = tabula.read_pdf(  # type: ignore[attr-defined]
                tables_source,
                pages=str(page_index),
                multiple_tables=True,
                pandas_options={"header": 0},
            )
            for idx, df in enumerate(tabula_tables, 1):
                caption_bits = []
                try:
                    caption_bits = [str(val).strip() for val in df.columns if str(val).strip()]
                except Exception:
                    caption_bits = []
                table_entries.append(
                    {
                        "id": f"Table {page_index}.{idx}",
                        "caption": "; ".join(caption_bits) or None,
                        "bbox": [0.0, 0.0, 0.0, 0.0],
                        "confidence": 0.55,
                    }
                )
        except Exception as tabula_exc:  # pragma: no cover - optional dependency
            logger.debug("tabula-py failed on page %s: %s", page_index, tabula_exc)
            tables_mode = "lines"

    if tables_mode == "lines":
        for caption in table_captions:
            table_entries.append(
                {
                    "id": f"Table {caption['number']}",
                    "caption": caption["title"] or caption["raw"],
                    "bbox": _round_bbox(caption["bbox"]),
                    "confidence": 0.55,
                }
            )

    return figure_entries, table_entries, orphan_captions


def _build_section_path(page: Dict[str, Any]) -> List[str]:
    path: List[str] = []
    chapter_title = page.get("chapter_title")
    chapter_label = page.get("chapter")
    alt_label = page.get("alt_chapter_label")
    topic_label = page.get("topic")
    topic_title = page.get("topic_title")

    if chapter_title:
        path.append(str(chapter_title))
    elif chapter_label:
        path.append(str(chapter_label))

    if alt_label and str(alt_label) not in path:
        path.append(str(alt_label))

    if topic_title:
        path.append(str(topic_title))
    elif topic_label:
        label_value = str(topic_label)
        if label_value not in path:
            path.append(label_value)

    return [segment for segment in path if segment]


def _ensure_page_schema(page: Dict[str, Any]) -> None:
    defaults_str = [
        "chapter",
        "chapter_title",
    "alt_chapter_label",
        "topic",
        "topic_title",
        "heading_type",
        "heading_line",
        "search_heading",
        "summary",
        "difficulty_level",
    ]
    for key in defaults_str:
        if page.get(key) is None:
            page[key] = "" if key != "difficulty_level" else "unknown"
        else:
            page[key] = str(page[key]) if page[key] is not None else ""

    error_value = page.get("error")
    page["error"] = str(error_value) if error_value not in (None, "") else None

    defaults_list = [
        "keywords",
        "qa_pairs",
        "figures",
        "tables",
        "learning_objectives",
        "key_terms",
        "footnotes",
        "references",
        "sidebars",
        "section_path",
        "answer_key_lines",
        "assessment_blocks",
    ]
    for key in defaults_list:
        if key not in page or page[key] is None:
            page[key] = []

    page.setdefault("has_content", False)
    page.setdefault("error", "")

    cw = page.get("context_window") or {}
    page["context_window"] = {
        "previous_topic": (cw.get("previous_topic") or ""),
        "current_topic": (cw.get("current_topic") or ""),
        "next_topic": (cw.get("next_topic") or ""),
    }


def _validate_monotonic_sequences(
    pages: Dict[str, Dict[str, Any]],
    diagnostics: Dict[str, Any],
) -> None:
    diagnostics.setdefault("out_of_range_headings", [])
    previous_chapter: Optional[int] = None
    previous_topic: Optional[Tuple[int, ...]] = None
    ordered_pages = sorted(pages.values(), key=lambda p: p.get("page_number") or 0)
    for page in ordered_pages:
        chapter_num = _extract_first_int(page.get("chapter"))
        if chapter_num is not None:
            if previous_chapter is not None and chapter_num < previous_chapter:
                diagnostics["out_of_range_headings"].append(
                    {
                        "page": page.get("page_number"),
                        "chapter": page.get("chapter"),
                        "reason": "chapter_regression",
                    }
                )
            previous_chapter = chapter_num
        topic_label = page.get("topic")
        if topic_label:
            topic_tuple = _topic_tuple(str(topic_label))
            if previous_topic and topic_tuple < previous_topic:
                diagnostics["out_of_range_headings"].append(
                    {
                        "page": page.get("page_number"),
                        "topic": topic_label,
                        "reason": "topic_regression",
                    }
                )
            previous_topic = topic_tuple or previous_topic


def _detect_heading_from_lines(
    lines: List[str],
    page_number: Optional[int],
    chapter_range: Optional[Tuple[int, int]],
    topic_range: Optional[Tuple[int, int]],
    toc_titles: List[str],
    toc_chapter_title: Optional[str],
    diagnostics: Dict[str, Any],
    options: PDFParseOptions,
) -> Dict[str, Any]:
    """Inspect page lines for chapter/topic headings with strict gating."""

    if topic_range is not None:
        return {"heading_type": "none"}

    diagnostics.setdefault("heading_candidates_rejected", defaultdict(int))
    result: Dict[str, Any] = {"heading_type": "none"}

    if not lines:
        return result

    max_idx = max(1, int(len(lines) * 0.25))
    toc_titles_norm = [
        _normalize_heading_text(title).lower() for title in toc_titles if title
    ]

    for idx, raw in enumerate(lines[:8]):
        if not raw or not raw.strip():
            continue

        normalized = _normalize_heading_text(raw)
        if not normalized:
            diagnostics["heading_candidates_rejected"]["empty"] += 1
            continue

        lowered = normalized.lower()
        if lowered.startswith(HEADING_BLACKLIST_PREFIXES):
            diagnostics["heading_candidates_rejected"]["blacklist"] += 1
            continue

        if idx > max_idx:
            diagnostics["heading_candidates_rejected"]["position"] += 1
            continue

        punctuation_ratio = sum(ch in ".,;:?!" for ch in normalized) / max(1, len(normalized))
        if len(normalized) > 120 or punctuation_ratio > 0.06:
            diagnostics["heading_candidates_rejected"]["structure"] += 1
            continue

        synonym_match = CHAP_SYNONYM_RE.match(normalized)
        chapter_label = None
        chapter_num_value: Optional[str] = None
        chapter_title_value: Optional[str] = None
        synonym_payload: Optional[Dict[str, Any]] = None

        if synonym_match:
            label = synonym_match.group("label").lower()
            number = synonym_match.group("num")
            title_fragment = synonym_match.group("title").strip() if synonym_match.group("title") else ""
            chapter_label = _format_synonym_label(label, number)
            chapter_num_value = number
            chapter_title_value = title_fragment or None
            if label in SYNONYM_CHAPTER_LABELS:
                if label == "lecture" and options.heading_lexicon not in {"auto", "lecture"}:
                    synonym_match = None
                    chapter_label = None
                    chapter_num_value = None
                    chapter_title_value = None
                else:
                    synonym_payload = {
                        "label": label,
                        "number": number,
                        "title": chapter_title_value,
                        "display": chapter_label,
                    }
        pattern_is_chapter = bool(synonym_match)

        legacy_chapter_match = None
        if not pattern_is_chapter:
            legacy_chapter_match = CHAPTER_LINE_RE.match(normalized)
            pattern_is_chapter = bool(legacy_chapter_match)
            if legacy_chapter_match:
                chapter_num_value = legacy_chapter_match.group(1)
                chapter_title_value = legacy_chapter_match.group(2).strip() if legacy_chapter_match.group(2) else None
                chapter_label = _format_synonym_label("chapter", chapter_num_value)

        topic_match = TOPIC_LINE_RE.match(normalized)
        if not topic_match:
            topic_match = TOPIC_NUM_RE.match(normalized)
        pattern_is_topic = bool(topic_match)

        fuzzy_hit = False
        if toc_titles_norm:
            for title in toc_titles_norm:
                if not title:
                    continue
                ratio = SequenceMatcher(None, lowered, title).ratio()
                if ratio >= 0.85:
                    fuzzy_hit = True
                    break

        prepared = normalized

        if not (pattern_is_chapter or pattern_is_topic or fuzzy_hit):
            lecture_enabled = options.heading_lexicon in {"auto", "lecture"}
            if lecture_enabled and re.match(r"^\s*lecture\s+\d", prepared, re.IGNORECASE):
                continue

            subhead_match = None
            for keyword in INLINE_TOPIC_KEYWORDS:
                if prepared.lower().startswith(keyword):
                    subhead_match = keyword
                    break
            if subhead_match:
                content_lines = _collect_subsection_lines(lines, idx)
                result_topic = {
                    "heading_type": "topic",
                    "topic": None,
                    "topic_title": prepared,
                    "heading_line": prepared,
                    "search_heading": _search_simplify(prepared),
                    "_topic_content": content_lines,
                }
                if re.match(r"^(how to|tips?|checklist)", prepared, re.IGNORECASE):
                    result_topic["_sidebar_entry"] = {
                        "title": prepared,
                        "content": content_lines,
                    }
                if re.match(r"^(suggested reading|further reading)", prepared, re.IGNORECASE):
                    result_topic["_reference_block"] = content_lines
                return result_topic

            diagnostics["heading_candidates_rejected"]["pattern"] += 1
            continue

        result["heading_type"] = "chapter" if pattern_is_chapter else "topic"
        result["heading_line"] = normalized
        result["search_heading"] = _search_simplify(normalized)

        if pattern_is_chapter and (chapter_label or chapter_num_value):
            if synonym_payload:
                resolved_title_syn = chapter_title_value or toc_chapter_title or normalized
                result["chapter_title"] = resolved_title_syn
                result["chapter"] = chapter_label or normalized
                result["alt_chapter_label"] = chapter_label
                result["_chapter_synonym"] = synonym_payload
                return result
            chapter_name = chapter_label or _format_synonym_label("chapter", chapter_num_value or "")
            result["chapter"] = chapter_name
            resolved_title = chapter_title_value or toc_chapter_title or normalized
            result["chapter_title"] = resolved_title
            return result

        topic_num = None
        topic_title = None
        if topic_match:
            topic_num = topic_match.group(1)
            topic_title = topic_match.group(2).strip() if topic_match.group(2) else ""
        if not topic_num:
            diagnostics["heading_candidates_rejected"]["topic_parse"] += 1
            return {"heading_type": "none"}

        result["topic"] = topic_num
        title_part = topic_title or normalized
        if normalized.lower().startswith(topic_num.lower()):
            title_part = normalized[len(topic_num) :].lstrip(" .-—") or title_part
        result["topic_title"] = title_part or normalized
        return result

    return result


# ========================= Helpers: text cleaning (PDF) =========================

def _clean_text(text: str) -> str:
    """Turn raw PDF text into a single readable paragraph string.

    - NFC normalization
    - Replace smart quotes/dashes; drop soft hyphen
    - Remove page/line numbers and punctuation-only lines
    - Fix hyphenation across line breaks
    - Merge lines into sentences/paragraphs
    - Heuristically drop short repeated headers/footers
    - Collapse whitespace
    """
    if not text:
        return ""

    s = unicodedata.normalize("NFC", text)
    s = s.replace("\u00A0", " ").replace("\u200B", "").replace("\ufeff", "")
    for k, v in {
        "\u2018": "'", "\u2019": "'", "\u201B": "'", "\u2032": "'",
        "\u201C": '"', "\u201D": '"', "\u2033": '"',
        "\u2013": "-", "\u2014": "-", "\u2212": "-", "\u00AD": "",
    }.items():
        s = s.replace(k, v)

    # Fix hyphenation across line breaks
    s = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", s)

    lines = s.splitlines()
    from collections import Counter
    stripped = [ln.strip() for ln in lines if ln.strip()]
    freq = Counter(stripped)
    header_footer = set()
    for cand in (stripped[:2] + stripped[-2:]):
        if cand and len(cand) < 40 and freq[cand] > 1:
            header_footer.add(cand)

    cleaned: List[str] = []
    for ln in lines:
        t = ln.strip()
        if not t:
            continue
        if t in header_footer:
            continue
        if re.fullmatch(r"\d+", t):
            continue
        if re.fullmatch(r"(?i)(?:page|p\.)\s*\d+(?:\s*of\s*\d+)?", t):
            continue
        if re.fullmatch(r"[\W_]+", t):
            continue
        cleaned.append(t)

    if not cleaned:
        return ""

    merged: List[str] = []
    for ln in cleaned:
        ln = re.sub(r"\s+", " ", ln).strip()
        if not ln:
            continue
        if not merged:
            merged.append(ln)
            continue
        if re.search(r"[\.!?]\"?\s*$", merged[-1]):
            merged.append(ln)
        else:
            merged[-1] = (merged[-1] + " " + ln).strip()

    out = " ".join(merged)
    out = re.sub(r"\s+", " ", out).strip()
    return out


# ========================= Post-processing helpers (PDF) =========================


def _parse_toc_structure(
    pages: Dict[str, Dict[str, Any]],
    total_pages: int,
    options: PDFParseOptions,
) -> Tuple[
    Dict[str, Dict[str, Any]],
    Dict[str, Tuple[int, int]],
    Dict[str, Tuple[int, int]],
    Dict[str, Dict[str, Any]],
    Optional[int],
]:
    """Parse Table of Contents pages, normalise entries, and compute ranges."""

    toc: Dict[str, Dict[str, Any]] = {}
    chapter_ranges: Dict[str, Tuple[int, int]] = {}
    topic_ranges: Dict[str, Tuple[int, int]] = {}
    synonym_index: Dict[str, Dict[str, Any]] = {}

    ordered_keys = sorted(pages.keys(), key=lambda k: pages[k].get("page_number", 0))
    first_pages = [
        pages[k]
        for k in ordered_keys
        if 0 < (pages[k].get("page_number") or 0) <= options.toc_window
    ]

    synonym_entries: List[Dict[str, Any]] = []
    topic_entries: List[Dict[str, Any]] = []

    awaiting_canonical_title: Optional[str] = None
    awaiting_synonym_index: Optional[int] = None
    current_canonical_parent: Optional[str] = None
    first_entry_start: Optional[int] = None

    for page in first_pages:
        lines = page.get("lines", []) or []
        for raw in lines:
            title_line, located_page = _split_toc_line(raw)
            if not title_line:
                continue

            title_line = re.sub(r"\.{2,}", " ", title_line)
            title_line = re.sub(r"\s+", " ", title_line).strip()

            normalized = _normalize_heading_text(title_line)
            if not normalized:
                continue

            lowered = normalized.lower()
            if lowered in {"contents", "table of contents"}:
                awaiting_canonical_title = None
                awaiting_synonym_index = None
                continue

            if located_page is None:
                if awaiting_canonical_title and awaiting_canonical_title in toc:
                    toc[awaiting_canonical_title]["chapter_title"] = normalized
                    awaiting_canonical_title = None
                    continue
                if (
                    awaiting_synonym_index is not None
                    and 0 <= awaiting_synonym_index < len(synonym_entries)
                ):
                    synonym_entries[awaiting_synonym_index]["title"] = normalized
                    awaiting_synonym_index = None
                continue

            synonym_match = CHAP_SYNONYM_RE.match(normalized)
            if synonym_match:
                label = synonym_match.group("label").lower()
                number = synonym_match.group("num")
                title_fragment = (synonym_match.group("title") or "").strip()
                display = _format_synonym_label(label, number)

                first_entry_start = located_page if first_entry_start is None else min(first_entry_start, located_page)

                if label in CANONICAL_CHAPTER_LABELS:
                    entry = toc.setdefault(number, {
                        "chapter_title": None,
                        "start_page": located_page,
                        "topics": [],
                        "label": label,
                    })
                    if entry.get("start_page") is None or located_page < entry.get("start_page", located_page):
                        entry["start_page"] = located_page
                    if title_fragment:
                        entry["chapter_title"] = title_fragment
                    else:
                        awaiting_canonical_title = number
                    current_canonical_parent = number
                else:
                    synonym_entries.append({
                        "label": label,
                        "number": number,
                        "display": display,
                        "title": title_fragment or None,
                        "start_page": located_page,
                        "parent": current_canonical_parent,
                    })
                    if not title_fragment:
                        awaiting_synonym_index = len(synonym_entries) - 1
                continue

            topic_match = TOPIC_LINE_RE.match(normalized) or TOPIC_NUM_RE.match(normalized)
            if topic_match:
                topic_num = _normalize_topic_number(topic_match.group(1)) if topic_match.group(1) else None
                if not topic_num:
                    continue
                topic_title = _normalize_heading_text(topic_match.group(2) or "") or None
                topic_entries.append({
                    "num": topic_num,
                    "title": topic_title,
                    "page": located_page,
                })
                first_entry_start = located_page if first_entry_start is None else min(first_entry_start, located_page)
                continue

    # Ensure canonical entries exist for topics that appeared without an explicit parent
    for topic_entry in topic_entries:
        chapter_key = topic_entry["num"].split(".")[0]
        entry = toc.setdefault(chapter_key, {
            "chapter_title": None,
            "start_page": topic_entry["page"],
            "topics": [],
            "label": "chapter",
        })
        if entry.get("start_page") is None or topic_entry["page"] < entry.get("start_page", topic_entry["page"]):
            entry["start_page"] = topic_entry["page"]
        entry.setdefault("topics", []).append({
            "num": topic_entry["num"],
            "title": topic_entry["title"],
            "page": topic_entry["page"],
        })

    # Derive chapter ranges based on start pages
    sorted_chapters = sorted(
        (
            entry.get("start_page"),
            key,
            entry,
        )
        for key, entry in toc.items()
        if entry.get("start_page")
    )

    for idx, (start, chapter_key, entry) in enumerate(sorted_chapters):
        if start is None:
            continue
        next_start = None
        if idx + 1 < len(sorted_chapters):
            next_start = sorted_chapters[idx + 1][0]
        end = (next_start - 1) if next_start and next_start > start else total_pages
        entry["end_page"] = end
        chapter_ranges[chapter_key] = (start, end)

        topics = sorted(entry.get("topics", []), key=lambda t: t.get("page") or start)
        for topic_idx, topic in enumerate(topics):
            topic_start = topic.get("page") or start
            next_topic_start = None
            if topic_idx + 1 < len(topics):
                next_topic_start = topics[topic_idx + 1].get("page")
            boundaries = [candidate for candidate in [next_topic_start, next_start, end + 1] if candidate and candidate > topic_start]
            boundary = min(boundaries) if boundaries else end + 1
            topic_end = boundary - 1 if boundary > topic_start else end
            topic["start_page"] = topic_start
            topic["end_page"] = max(topic_start, topic_end)
            topic_ranges[topic["num"]] = (topic["start_page"], topic["end_page"])

    # Compute synonym ranges relative to canonical parents or sequential order
    synonym_entries.sort(key=lambda item: item["start_page"])
    for idx, entry in enumerate(synonym_entries):
        next_start = synonym_entries[idx + 1]["start_page"] if idx + 1 < len(synonym_entries) else None
        parent_range = chapter_ranges.get(entry["parent"] or "") if entry.get("parent") else None
        end_candidates: List[int] = []
        if next_start and next_start > entry["start_page"]:
            end_candidates.append(next_start - 1)
        if parent_range:
            end_candidates.append(parent_range[1])
        if not end_candidates:
            end_candidates.append(total_pages)
        end_page = max(entry["start_page"], min(end_candidates))
        entry["end_page"] = end_page
        key = f"{entry['label']} {entry['number']}".lower()
        synonym_index[key] = entry

    if first_entry_start is None and synonym_entries:
        first_entry_start = min(entry["start_page"] for entry in synonym_entries)

    return toc, chapter_ranges, topic_ranges, synonym_index, first_entry_start


def _map_synonym_to_chapter(
    page_number: int,
    synonym_data: Dict[str, Any],
    synonym_index: Dict[str, Dict[str, Any]],
    chapter_ranges: Dict[str, Tuple[int, int]],
    toc: Dict[str, Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    label_key = f"{synonym_data.get('label', '')} {synonym_data.get('number', '')}".lower().strip()
    entry = synonym_index.get(label_key)
    if not entry:
        return None

    chapter_key: Optional[str] = None
    reason: Optional[str] = None

    parent_hint = entry.get("parent")
    if parent_hint and parent_hint in chapter_ranges:
        chapter_key = parent_hint
        reason = "toc_parent"

    if chapter_key is None:
        for candidate_key, bounds in chapter_ranges.items():
            start, end = bounds
            if start <= entry["start_page"] <= end:
                chapter_key = candidate_key
                reason = "toc_range"
                break

    if chapter_key is None:
        sorted_ranges = sorted((bounds[0], key, bounds) for key, bounds in chapter_ranges.items())
        previous: Optional[Tuple[str, Tuple[int, int]]] = None
        for start, key, bounds in sorted_ranges:
            if entry["start_page"] >= start:
                previous = (key, bounds)
            else:
                break
        if previous:
            chapter_key = previous[0]
            reason = "nearest_previous"

    if not chapter_key:
        return None

    chapter_entry = toc.get(chapter_key, {})
    chapter_label = chapter_entry.get("label") or "chapter"
    canonical_label = _format_synonym_label(chapter_label, chapter_key)
    chapter_title = chapter_entry.get("chapter_title") or canonical_label

    mapped_info = {
        "chapter": canonical_label,
        "chapter_title": chapter_title,
        "reason": reason or "inferred",
        "chapter_key": chapter_key,
        "entry": entry,
    }

    if chapter_key in chapter_ranges:
        start, end = chapter_ranges[chapter_key]
        if not (start <= page_number <= end):
            mapped_info["out_of_range"] = True

    return mapped_info


def _build_topic_reverse_index(
    toc: Dict[str, Dict[str, Any]],
    chapter_ranges: Dict[str, Tuple[int, int]],
    topic_ranges: Dict[str, Tuple[int, int]],
) -> Dict[str, Dict[str, Any]]:
    reverse: Dict[str, Dict[str, Any]] = {}
    for chapter_num, info in toc.items():
        chapter_title = info.get("chapter_title")
        for topic in info.get("topics", []):
            num = topic.get("num")
            if not num:
                continue
            start_page, end_page = topic_ranges.get(num, (topic.get("start_page"), topic.get("end_page")))
            reverse[num] = {
                "chapter": chapter_num,
                "title": topic.get("title"),
                "start_page": start_page,
                "end_page": end_page,
                "chapter_title": chapter_title,
            }
    return reverse


def _postprocess_headings_and_context(
    pages: Dict[str, Dict[str, Any]],
    toc: Dict[str, Dict[str, Any]],
    chapter_ranges: Dict[str, Tuple[int, int]],
    topic_ranges: Dict[str, Tuple[int, int]],
    topic_index: Dict[str, Dict[str, Any]],
    synonym_index: Dict[str, Dict[str, Any]],
    first_entry_start: Optional[int],
    diagnostics: Dict[str, Any],
    options: PDFParseOptions,
) -> None:
    """Propagate headings strictly within TOC bounds and fill context windows."""

    order = sorted(pages.keys(), key=lambda k: pages[k].get("page_number", 0))
    diagnostics.setdefault("synonym_mappings", [])
    diagnostics.setdefault("unmapped_headings", [])

    diagnostics.setdefault("no_toc_lecture_mode_used", False)
    diagnostics.setdefault("synthesized_ranges", [])
    if not chapter_ranges and synonym_index:
        sorted_synonyms = sorted(
            synonym_index.values(), key=lambda entry: entry.get("start_page") or 0
        )
        synthesized: List[Dict[str, Any]] = []
        for idx, entry in enumerate(sorted_synonyms):
            start = entry.get("start_page") or 0
            if start <= 0:
                continue
            next_start = (
                sorted_synonyms[idx + 1].get("start_page")
                if idx + 1 < len(sorted_synonyms)
                else None
            )
            candidate_end = entry.get("end_page") or start
            if next_start and next_start > start:
                candidate_end = min(candidate_end, next_start - 1)
            end = max(start, candidate_end)
            label_display = entry.get("display") or entry.get("label") or ""
            synthesized.append({"label": label_display, "start": start, "end": end})
            number_key = entry.get("number")
            if number_key and number_key not in chapter_ranges:
                chapter_ranges[number_key] = (start, end)
                toc.setdefault(
                    number_key,
                    {
                        "chapter_title": entry.get("title") or label_display,
                        "start_page": start,
                        "end_page": end,
                        "topics": [],
                        "label": entry.get("label"),
                    },
                )
            for key_inner in order:
                page = pages[key_inner]
                pn = page.get("page_number") or 0
                if start <= pn <= end:
                    page.setdefault("chapter", label_display)
                    if entry.get("title"):
                        page.setdefault("chapter_title", entry["title"])
                    page.setdefault("alt_chapter_label", label_display)
                    if not page.get("heading_type") or page.get("heading_type") == "none":
                        page["heading_type"] = "chapter"
        if synthesized:
            diagnostics["no_toc_lecture_mode_used"] = True
            diagnostics["synthesized_ranges"] = synthesized

    chapter_sequence: List[Tuple[int, int, str, Optional[str], str]] = []
    for chapter_num, bounds in chapter_ranges.items():
        toc_entry = toc.get(chapter_num, {})
        chapter_title = toc_entry.get("chapter_title")
        chapter_label = toc_entry.get("label") or "chapter"
        chapter_display = chapter_title or _format_synonym_label(chapter_label, chapter_num)
        chapter_sequence.append((bounds[0], bounds[1], chapter_num, chapter_title, chapter_display))
    chapter_sequence.sort(key=lambda item: item[0])

    topic_sequence: List[Tuple[int, int, str, Optional[str], Optional[str]]] = []
    for topic_num, bounds in topic_ranges.items():
        info = topic_index.get(topic_num, {})
        topic_sequence.append(
            (bounds[0], bounds[1], topic_num, info.get("title"), info.get("chapter"))
        )
    topic_sequence.sort(key=lambda item: item[0])

    chapter_ptr = 0
    topic_ptr = 0
    active_chapter: Optional[str] = None
    active_chapter_title: Optional[str] = None
    active_chapter_end: Optional[int] = None
    active_topic: Optional[str] = None
    active_topic_title: Optional[str] = None
    active_topic_end: Optional[int] = None

    diagnostics.setdefault("front_matter_pages", [])

    # Collect topic labels for context windows later
    for key in order:
        page = pages[key]
        pn = page.get("page_number")
        if pn is None:
            continue

        if first_entry_start is not None and pn < first_entry_start:
            page["heading_type"] = "front_matter"
            page["chapter"] = None
            page["chapter_title"] = None
            page["topic"] = None
            page["topic_title"] = None
            if page.get("has_content"):
                diagnostics["front_matter_pages"].append(pn)
            continue

        while chapter_ptr < len(chapter_sequence) and pn > chapter_sequence[chapter_ptr][1]:
            chapter_ptr += 1

        chapter_num: Optional[str] = None
        chapter_title: Optional[str] = None
        chapter_end: Optional[int] = None

        if chapter_ptr < len(chapter_sequence):
            start, end, current_chapter_num, current_title, current_display = chapter_sequence[chapter_ptr]
            if start <= pn <= end:
                chapter_num = current_chapter_num
                chapter_title = current_title
                chapter_end = end

        if chapter_num:
            toc_entry = toc.get(chapter_num, {})
            label_value = toc_entry.get("label") or "chapter"
            active_chapter = _format_synonym_label(label_value, chapter_num)
            active_chapter_title = chapter_title
            active_chapter_end = chapter_end
            page["chapter"] = active_chapter
            if chapter_title:
                page["chapter_title"] = chapter_title
        elif active_chapter and active_chapter_end and pn <= active_chapter_end:
            page.setdefault("chapter", active_chapter)
            if active_chapter_title:
                page.setdefault("chapter_title", active_chapter_title)
        else:
            active_chapter = None
            active_chapter_title = None
            active_chapter_end = None

        while topic_ptr < len(topic_sequence) and pn > topic_sequence[topic_ptr][1]:
            topic_ptr += 1

        topic_num: Optional[str] = None
        topic_title: Optional[str] = None
        topic_end: Optional[int] = None
        topic_chapter: Optional[str] = None

        if topic_ptr < len(topic_sequence):
            start, end, current_topic_num, current_topic_title, topic_chapter_num = topic_sequence[topic_ptr]
            if start <= pn <= end:
                topic_num = current_topic_num
                topic_title = current_topic_title
                topic_end = end
                topic_chapter = topic_chapter_num

        if topic_num:
            page["topic"] = topic_num
            if topic_title:
                page["topic_title"] = topic_title
            if topic_chapter:
                page.setdefault("chapter", f"Chapter {topic_chapter}")
                chapter_title = toc.get(str(topic_chapter), {}).get("chapter_title")
                if chapter_title:
                    page.setdefault("chapter_title", chapter_title)
            active_topic = topic_num
            active_topic_title = topic_title
            active_topic_end = topic_end
        else:
            if active_topic and active_topic_end and pn <= active_topic_end:
                page.setdefault("topic", active_topic)
                if active_topic_title:
                    page.setdefault("topic_title", active_topic_title)
            else:
                active_topic = None
                active_topic_title = None
                active_topic_end = None

        chapter_range = None
        chapter_num_for_lookup: Optional[str] = None
        if page.get("chapter"):
            chapter_num_value = _extract_first_int(page.get("chapter"))
            if chapter_num_value is not None:
                chapter_num_for_lookup = str(chapter_num_value)
                chapter_range = chapter_ranges.get(chapter_num_for_lookup)

        topic_range = topic_ranges.get(page.get("topic")) if page.get("topic") else None

        if not page.get("topic"):
            toc_titles: List[str] = []
            if chapter_num_for_lookup and chapter_num_for_lookup in toc:
                chapter_info = toc[chapter_num_for_lookup]
                if chapter_info.get("chapter_title"):
                    toc_titles.append(chapter_info["chapter_title"])
                toc_titles.extend(
                    t.get("title")
                    for t in chapter_info.get("topics", [])
                    if t.get("title")
                )

            detection = _detect_heading_from_lines(
                page.get("lines", []) or [],
                pn,
                chapter_range,
                topic_range,
                toc_titles,
                toc.get(chapter_num_for_lookup, {}).get("chapter_title") if chapter_num_for_lookup else None,
                diagnostics,
                options,
            )

            if detection.get("heading_type") and detection.get("heading_type") != "none":
                page["heading_line"] = detection.get("heading_line", "")
                page["search_heading"] = detection.get("search_heading", "")

            if detection.get("heading_type") == "topic":
                candidate_topic = detection.get("topic")
                if candidate_topic and candidate_topic in topic_ranges:
                    start, end = topic_ranges[candidate_topic]
                    if start <= pn <= end:
                        page["topic"] = candidate_topic
                        page["topic_title"] = detection.get("topic_title")
                        active_topic = candidate_topic
                        active_topic_title = detection.get("topic_title")
                        active_topic_end = end
                        info = topic_index.get(candidate_topic, {})
                        chapter_num_for_lookup = info.get("chapter") or chapter_num_for_lookup
                        if info.get("chapter"):
                            page.setdefault("chapter", f"Chapter {info['chapter']}")
                            chapter_data = toc.get(str(info.get("chapter")), {})
                            if chapter_data.get("chapter_title"):
                                page.setdefault("chapter_title", chapter_data["chapter_title"])
                else:
                    diagnostics.setdefault("mismatch_pages", []).append({
                        "page": pn,
                        "assigned": detection.get("topic"),
                        "expected": None,
                        "reason": "topic_not_in_toc",
                    })

                sidebar_entry = detection.get("_sidebar_entry")
                if sidebar_entry and sidebar_entry.get("content"):
                    page.setdefault("sidebars", []).append(sidebar_entry)
                reference_block = detection.get("_reference_block")
                if reference_block:
                    page.setdefault("references", [])
                    for ref_line in reference_block:
                        if ref_line and ref_line not in page["references"]:
                            page["references"].append(ref_line)

            elif detection.get("heading_type") == "chapter" and not page.get("chapter"):
                candidate_chapter_num = _extract_first_int(detection.get("chapter"))
                if candidate_chapter_num is not None:
                    candidate_key = str(candidate_chapter_num)
                    if candidate_key in chapter_ranges:
                        start, end = chapter_ranges[candidate_key]
                        if start <= pn <= end:
                            parent_entry = toc.get(candidate_key, {})
                            parent_label = parent_entry.get("label") or "chapter"
                            page["chapter"] = _format_synonym_label(parent_label, candidate_key)
                            page["chapter_title"] = detection.get("chapter_title") or parent_entry.get("chapter_title")
                            active_chapter = page["chapter"]
                            active_chapter_title = page.get("chapter_title")
                            active_chapter_end = end

            synonym_payload = detection.get("_chapter_synonym") if detection else None
            if synonym_payload:
                alt_display = synonym_payload.get("display") or detection.get("chapter") if detection else None
                if alt_display:
                    page["alt_chapter_label"] = alt_display
                if options.heading_synonyms:
                    mapped = _map_synonym_to_chapter(
                        pn,
                        synonym_payload,
                        synonym_index,
                        chapter_ranges,
                        toc,
                    )
                    if mapped:
                        page["chapter"] = mapped["chapter"]
                        page["chapter_title"] = mapped.get("chapter_title") or page.get("chapter_title")
                        diagnostics.setdefault("synonym_mappings", []).append({
                            "page": pn,
                            "from": alt_display,
                            "to": mapped["chapter"],
                            "reason": mapped.get("reason"),
                        })
                        chapter_num_for_lookup = mapped.get("chapter_key") or chapter_num_for_lookup
                        if chapter_num_for_lookup and chapter_num_for_lookup in chapter_ranges:
                            active_chapter = mapped["chapter"]
                            active_chapter_title = page.get("chapter_title")
                            active_chapter_end = chapter_ranges[chapter_num_for_lookup][1]
                    else:
                        diagnostics.setdefault("unmapped_headings", []).append({
                            "page": pn,
                            "label": alt_display,
                        })
                        page.setdefault("chapter", alt_display)
                        page.setdefault("chapter_title", detection.get("chapter_title") if detection else alt_display)
                else:
                    diagnostics.setdefault("unmapped_headings", []).append({
                        "page": pn,
                        "label": alt_display,
                        "reason": "synonyms_disabled",
                    })
                    page.setdefault("chapter", alt_display)
                    page.setdefault("chapter_title", detection.get("chapter_title") if detection else alt_display)

        if not page.get("alt_chapter_label") and page.get("chapter"):
            chapter_lower = str(page["chapter"]).lower()
            for label in SYNONYM_CHAPTER_LABELS:
                if chapter_lower.startswith(label):
                    page["alt_chapter_label"] = page["chapter"]
                    break

        if page.get("topic"):
            page["heading_type"] = "topic"
        elif page.get("chapter"):
            page["heading_type"] = "chapter"
        else:
            page["heading_type"] = "none"

    mismatch_count = len(diagnostics.get("mismatch_pages", []))
    if mismatch_count:
        logger.warning(
            "Validation: %d heading assignments cleared (outside TOC bounds)",
            mismatch_count,
        )

    # Build context windows and section paths using topics when available, fallback to chapters/lectures otherwise
    topic_sequence_sorted = sorted(topic_sequence, key=lambda item: item[0])
    topic_nav: List[Tuple[int, int, str]] = [
        (start, end, title or num)
        for start, end, num, title, _ in topic_sequence_sorted
    ]

    chapter_nav: List[Tuple[int, int, str]] = [
        (start, end, display)
        for start, end, _, _, display in chapter_sequence
    ]

    synonym_nav: List[Tuple[int, int, str]] = sorted(
        [
            (
                entry.get("start_page") or 0,
                entry.get("end_page") or entry.get("start_page") or 0,
                entry.get("display") or entry.get("label") or "",
            )
            for entry in synonym_index.values()
            if entry.get("start_page")
        ],
        key=lambda item: item[0],
    )

    if not chapter_nav and synonym_nav:
        chapter_nav = synonym_nav[:]

    fallback_nav: List[Tuple[int, int, str]] = []
    if not topic_nav and not chapter_nav and not synonym_nav:
        seen_label: Optional[str] = None
        for key in order:
            page = pages[key]
            pn = page.get("page_number") or 0
            label = page.get("chapter_title") or page.get("alt_chapter_label") or page.get("chapter")
            if not label:
                continue
            if label != seen_label:
                fallback_nav.append((pn, pn, label))
                seen_label = label

    for key in order:
        page = pages[key]
        pn = page.get("page_number") or 0
        current_label = (
            page.get("topic_title")
            or page.get("topic")
            or page.get("chapter_title")
            or page.get("alt_chapter_label")
            or page.get("chapter")
        )

        if page.get("topic") and topic_nav:
            nav_source = topic_nav
        elif page.get("alt_chapter_label") and synonym_nav:
            nav_source = synonym_nav
        elif chapter_nav:
            nav_source = chapter_nav
        elif fallback_nav:
            nav_source = fallback_nav
        else:
            nav_source = []

        prev_label = ""
        next_label = ""
        for start, end, label in nav_source:
            if end < start:
                end = start
            if pn < start:
                next_label = label
                break
            if start <= pn <= end:
                if not current_label:
                    current_label = label
            if start <= pn:
                prev_label = label

        cw = page.get("context_window") or {
            "previous_topic": "",
            "current_topic": "",
            "next_topic": "",
        }
        cw["previous_topic"] = prev_label
        cw["current_topic"] = current_label or ""
        cw["next_topic"] = next_label
        page["context_window"] = cw
        page["section_path"] = _build_section_path(page)

    last_path: List[str] = []
    for key in order:
        page = pages[key]
        current_path = page.get("section_path") or []
        if current_path:
            last_path = list(current_path)
        elif last_path:
            page["section_path"] = list(last_path)


def _validate_heading_ranges(
    pages: Dict[str, Dict[str, Any]],
    toc: Dict[str, Dict[str, Any]],
    chapter_ranges: Dict[str, Tuple[int, int]],
    topic_ranges: Dict[str, Tuple[int, int]],
    topic_index: Dict[str, Dict[str, Any]],
    diagnostics: Dict[str, Any],
) -> None:
    order = sorted(pages.keys(), key=lambda k: pages[k].get("page_number", 0))

    diagnostics.setdefault("mismatch_pages", [])
    diagnostics.setdefault("out_of_range_headings", [])

    for key in order:
        page = pages[key]
        pn = page.get("page_number")
        if pn is None:
            continue

        chapter = page.get("chapter")
        topic = page.get("topic")

        if chapter:
            chapter_num_value = _extract_first_int(chapter)
            chapter_key = str(chapter_num_value) if chapter_num_value is not None else None
            if (
                chapter_key is None
                or chapter_key not in chapter_ranges
                or not (chapter_ranges[chapter_key][0] <= pn <= chapter_ranges[chapter_key][1])
            ):
                diagnostics["mismatch_pages"].append({
                    "page": pn,
                    "assigned": chapter,
                    "expected": None,
                    "reason": "chapter_outside_toc_range",
                })
                diagnostics["out_of_range_headings"].append({
                    "page": pn,
                    "heading": chapter,
                    "type": "chapter",
                })
                page["chapter"] = None
                page["chapter_title"] = None
                page.pop("alt_chapter_label", None)
                if page.get("heading_type") == "chapter" and not topic:
                    page["heading_type"] = "none"
                chapter = None

        if topic:
            if topic not in topic_ranges or not (topic_ranges[topic][0] <= pn <= topic_ranges[topic][1]):
                repaired = _auto_repair_topic_assignment(
                    page,
                    pn,
                    toc,
                    chapter_ranges,
                    topic_ranges,
                    diagnostics,
                    reason="outside_toc_range",
                )
                if not repaired:
                    diagnostics["mismatch_pages"].append({
                        "page": pn,
                        "assigned": topic,
                        "expected": None,
                        "reason": "topic_outside_toc_range",
                    })
                    diagnostics["out_of_range_headings"].append({
                        "page": pn,
                        "heading": topic,
                        "type": "topic",
                    })
                    page["topic"] = None
                    page["topic_title"] = None
                    if page.get("heading_type") == "topic":
                        page["heading_type"] = "chapter" if page.get("chapter") else "none"
                    topic = None
            else:
                expected_chapter = topic_index.get(topic, {}).get("chapter")
                if expected_chapter and expected_chapter not in (None, ""):
                    toc_entry = toc.get(str(expected_chapter), {})
                    parent_label = toc_entry.get("label") or "chapter"
                    expected_label = _format_synonym_label(parent_label, str(expected_chapter))
                    if page.get("chapter") != expected_label:
                        page["chapter"] = expected_label
                        chapter_title = toc_entry.get("chapter_title")
                        if chapter_title:
                            page["chapter_title"] = chapter_title

        if page.get("topic"):
            page["heading_type"] = "topic"
        elif page.get("chapter"):
            page["heading_type"] = "chapter"
        else:
            page["heading_type"] = "none"


# ========================= Parser implementation =========================

class DocumentParser:
    """Service for parsing various document types and extracting content."""

    @staticmethod
    async def parse_document(
        file_content: bytes,
        document_type: DocumentType,
        filename: str,
        options: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        try:
            logger.info(f"Parsing {document_type.value} document: {filename}")
            if document_type == DocumentType.PDF:
                return await DocumentParser._parse_pdf(file_content, filename, options)
            if document_type == DocumentType.PPTX:
                return await DocumentParser._parse_pptx(file_content, filename)
            if document_type == DocumentType.DOCX:
                return await DocumentParser._parse_docx(file_content, filename)
            raise ValueError(f"Unsupported document type: {document_type}")
        except Exception as e:
            logger.error(f"Error parsing document {filename}: {str(e)}")
            raise

    @staticmethod
    async def _parse_pdf(
        file_content: bytes,
        filename: str,
        options: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        pdf_options = PDFParseOptions.from_dict(options)
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        pdf_metadata = pdf_reader.metadata or {}
        metadata = {
            "filename": filename,
            "file_type": "PDF",
            "parsed_at": datetime.utcnow().isoformat(),
            "document_type": "pdf",
            "parser_version": "v2.1",
            "date_parsed": datetime.utcnow().isoformat(),
            "total_pages": len(pdf_reader.pages),
            "title": pdf_metadata.get("/Title", "") or "",
            "author": pdf_metadata.get("/Author", "") or "",
            "subject": pdf_metadata.get("/Subject", "") or "",
            "creator": pdf_metadata.get("/Creator", "") or "",
            "producer": pdf_metadata.get("/Producer", "") or "",
            "creation_date": str(pdf_metadata.get("/CreationDate", "")) or "",
            "modification_date": str(pdf_metadata.get("/ModDate", "")) or "",
        }

        pages_content: Dict[str, Any] = {}
        total_words = 0
        pages_with_content = 0
        pages_with_errors = 0
        ocr_histogram: Dict[str, int] = {"high": 0, "medium": 0, "low": 0}
        figure_caption_orphans: List[Dict[str, Any]] = []
        ocr_trigger_threshold = OCR_THRESHOLDS.get(pdf_options.ocr_threshold, 60)
        temp_pdf_path: Optional[str] = None

        with ExitStack() as stack:
            doc_fitz = None
            if fitz is not None:
                try:
                    doc_fitz = fitz.open(stream=file_content, filetype="pdf")  # type: ignore[attr-defined]
                    stack.callback(doc_fitz.close)
                except Exception as exc:  # pragma: no cover - optional dependency
                    logger.warning("PyMuPDF unavailable for figure detection: %s", exc)
                    doc_fitz = None

            plumber_doc = None
            if pdfplumber is not None:
                try:
                    plumber_doc = pdfplumber.open(io.BytesIO(file_content))
                    stack.enter_context(plumber_doc)
                except Exception as exc:  # pragma: no cover - optional dependency
                    logger.debug("pdfplumber fallback disabled: %s", exc)
                    plumber_doc = None

            if pdf_options.tables_mode in {"camelot", "tabula"}:
                dependency_ready = (
                    (pdf_options.tables_mode == "camelot" and camelot is not None)
                    or (pdf_options.tables_mode == "tabula" and tabula is not None)
                )
                if dependency_ready:
                    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                    try:
                        tmp_file.write(file_content)
                        tmp_file.flush()
                        temp_pdf_path = tmp_file.name
                    finally:
                        tmp_file.close()

                    def _cleanup_temp(path: str) -> None:
                        try:
                            if os.path.exists(path):
                                os.remove(path)
                        except Exception:
                            pass

                    stack.callback(lambda path=temp_pdf_path: _cleanup_temp(path))

            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_key = f"Page {page_num}"
                page_content: Dict[str, Any] = {
                    "page_number": page_num,
                    "text": "",
                    "word_count": 0,
                    "char_count": 0,
                    "lines": [],
                    "has_content": False,
                    "error": None,
                    "chapter": None,
                    "chapter_title": None,
                    "alt_chapter_label": None,
                    "topic": None,
                    "topic_title": None,
                    "heading_type": "none",
                    "heading_line": "",
                    "search_heading": "",
                    "summary": "",
                    "keywords": [],
                    "qa_pairs": [],
                    "figures": [],
                    "tables": [],
                    "learning_objectives": [],
                    "key_terms": [],
                    "footnotes": [],
                    "references": [],
                    "sidebars": [],
                    "section_path": [],
                    "difficulty_level": "unknown",
                    "context_window": {
                        "previous_topic": "",
                        "current_topic": "",
                        "next_topic": "",
                    },
                    "page_source": page_num,
                    "content_cleaned": "",
                }

                fitz_page = None
                if doc_fitz is not None and 0 <= page_num - 1 < len(doc_fitz):
                    fitz_page = doc_fitz[page_num - 1]
                plumber_page = None
                if plumber_doc is not None and 0 <= page_num - 1 < len(plumber_doc.pages):
                    plumber_page = plumber_doc.pages[page_num - 1]

                try:
                    raw_text = page.extract_text() or ""
                    if (not raw_text.strip()) and plumber_page is not None:
                        try:
                            raw_text = plumber_page.extract_text() or ""
                        except Exception as fallback_exc:  # pragma: no cover - best effort
                            logger.debug("pdfplumber failed on page %s: %s", page_num, fallback_exc)

                    initial_words = len(raw_text.split()) if raw_text else 0
                    if fitz_page is not None and initial_words < ocr_trigger_threshold:
                        ocr_text = _run_ocr_on_page(fitz_page)
                        if ocr_text and len(ocr_text.split()) > initial_words:
                            raw_text = ocr_text
                        elif ocr_text:
                            raw_text = (raw_text + "\n" + ocr_text).strip()

                    if raw_text and raw_text.strip():
                        lines = [line.strip() for line in raw_text.split("\n") if line.strip()]
                        page_content["text"] = raw_text.strip()
                        page_content["word_count"] = len(raw_text.split())
                        page_content["char_count"] = len(raw_text.strip())
                        page_content["lines"] = lines
                        page_content["line_count"] = len(lines)
                        page_content["has_content"] = True

                        cleaned = _clean_text(page_content["text"])
                        if not cleaned and page_content["lines"]:
                            cleaned = " ".join(page_content["lines"])
                        page_content["content_cleaned"] = cleaned or page_content["text"]

                        page_content["learning_objectives"] = _extract_learning_objectives(lines)
                        page_content["key_terms"] = _extract_key_terms(lines)
                        footnotes, references = _extract_footnotes(lines)
                        page_content["footnotes"] = footnotes
                        page_content["references"] = references
                        page_content["sidebars"] = _extract_sidebars(lines)

                        total_words += page_content["word_count"]
                        pages_with_content += 1

                    try:
                        if hasattr(page, "mediabox"):
                            page_content["dimensions"] = {
                                "width": float(page.mediabox.width) if page.mediabox.width else 0,
                                "height": float(page.mediabox.height) if page.mediabox.height else 0,
                            }
                        if hasattr(page, "rotate"):
                            page_content["rotation"] = page.get("/Rotate", 0)
                    except Exception as meta_error:  # pragma: no cover - metadata best effort
                        logger.debug("Could not extract page %s metadata: %s", page_num, meta_error)

                except Exception as exc:
                    logger.warning("Error extracting text from page %s: %s", page_num, exc)
                    page_content["error"] = str(exc)
                    page_content["has_content"] = False
                    pages_with_errors += 1

                # Visual element detection (figures/tables)
                layout = None
                if fitz_page is not None:
                    try:
                        layout = fitz_page.get_text("rawdict")
                    except Exception as layout_exc:  # pragma: no cover - optional
                        logger.debug("PyMuPDF layout unavailable on page %s: %s", page_num, layout_exc)
                        layout = None

                figures, tables, orphans = _detect_visual_elements(
                    page_num,
                    fitz_page,
                    layout,
                    pdf_options,
                    tables_source=temp_pdf_path,
                )
                if figures:
                    page_content["figures"] = figures
                if tables:
                    page_content["tables"] = tables
                if orphans:
                    figure_caption_orphans.extend(orphans)

                histogram_bucket = _categorize_word_count(page_content.get("word_count", 0))
                ocr_histogram[histogram_bucket] = ocr_histogram.get(histogram_bucket, 0) + 1

                pages_content[page_key] = page_content

        toc, chapter_ranges, topic_ranges, synonym_index, first_entry_start = _parse_toc_structure(
            pages_content,
            len(pdf_reader.pages),
            pdf_options,
        )
        topic_index = _build_topic_reverse_index(toc, chapter_ranges, topic_ranges)
        diagnostics: Dict[str, Any] = {
            "heading_candidates_rejected": defaultdict(int),
        }
        _postprocess_headings_and_context(
            pages_content,
            toc,
            chapter_ranges,
            topic_ranges,
            topic_index,
            synonym_index,
            first_entry_start,
            diagnostics,
            pdf_options,
        )
        _validate_heading_ranges(
            pages_content,
            toc,
            chapter_ranges,
            topic_ranges,
            topic_index,
            diagnostics,
        )
        _validate_monotonic_sequences(pages_content, diagnostics)

        assessments, answer_keys = _aggregate_assessment_data(
            pages_content,
            pdf_options,
            diagnostics,
        )

        diagnostics.setdefault("front_matter_pages", [])
        diagnostics.setdefault("mismatch_pages", [])
        diagnostics.setdefault("auto_repairs", [])
        diagnostics.setdefault("heading_candidates_rejected", {})
        diagnostics.setdefault("out_of_range_headings", [])
        diagnostics.setdefault("no_toc_lecture_mode_used", False)
        diagnostics.setdefault("synthesized_ranges", [])
        diagnostics.setdefault("synonym_mappings", [])
        diagnostics.setdefault("unmapped_headings", [])
        diagnostics.setdefault("answer_key_pages_found", [])
        diagnostics.setdefault("assessment_blocks_found", [])
        diagnostics.setdefault("orphan_answers", [])
        diagnostics.setdefault("unmatched_questions", [])
        diagnostics["heading_candidates_rejected"] = dict(diagnostics.get("heading_candidates_rejected", {}))
        diagnostics["figure_caption_orphans"] = figure_caption_orphans
        diagnostics["ocr_confidence_histogram"] = ocr_histogram

        for page in pages_content.values():
            _ensure_page_schema(page)

        total_chapters_detected = len({
            page.get("chapter")
            for page in pages_content.values()
            if page.get("chapter")
        })
        total_topics_detected = len({
            page.get("topic")
            for page in pages_content.values()
            if page.get("topic")
        })

        logger.info(
            "Parse summary: total_chapters=%s total_topics=%s assessments_found=%s answer_keys_found=%s",
            total_chapters_detected,
            total_topics_detected,
            len(assessments),
            len(answer_keys),
        )

        result = {
            "type": "pdf",
            "metadata": metadata,
            "content": pages_content,
            "assessments": assessments,
            "answer_keys": answer_keys,
            "summary": {
                "total_pages": len(pdf_reader.pages),
                "pages_with_content": pages_with_content,
                "pages_with_errors": pages_with_errors,
                "empty_pages": len(pdf_reader.pages) - pages_with_content - pages_with_errors,
                "total_lines": sum(p.get("line_count", 0) for p in pages_content.values()),
                "total_characters": sum(p.get("char_count", 0) for p in pages_content.values()),
                "average_words_per_page": round(total_words / pages_with_content, 2) if pages_with_content > 0 else 0,
            },
            "total_word_count": total_words,
            "diagnostics": diagnostics,
        }

        return result

    @staticmethod
    async def _parse_pptx(file_content: bytes, filename: str) -> Dict[str, Any]:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)

        metadata = {
            "filename": filename,
            "file_type": "PPTX",
            "parsed_at": datetime.utcnow().isoformat(),
            "total_slides": len(presentation.slides),
            "title": presentation.core_properties.title or "",
            "author": presentation.core_properties.author or "",
            "subject": presentation.core_properties.subject or "",
            "created": presentation.core_properties.created.isoformat() if presentation.core_properties.created else "",
            "modified": presentation.core_properties.modified.isoformat() if presentation.core_properties.modified else "",
        }

        pages_content: Dict[str, Any] = {}

        for slide_num, slide in enumerate(presentation.slides, 1):
            page_key = f"Page {slide_num}"
            page_content = {
                "slide_number": slide_num,
                "title": "",
                "text_content": [],
                "images": [],
                "tables": [],
                "notes": "",
                "raw_content": [],
            }

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if not page_content["title"]:
                        page_content["title"] = shape.text.strip()
                    break

            for shape in slide.shapes:
                shape_data = DocumentParser._extract_shape_content(shape)
                if shape_data:
                    page_content["raw_content"].append(shape_data)
                    if shape_data["type"] == "text" and shape_data["text"].strip():
                        page_content["text_content"].append(shape_data["text"].strip())
                    if shape_data["type"] == "image":
                        page_content["images"].append(shape_data)
                    elif shape_data["type"] == "table":
                        page_content["tables"].append(shape_data)

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    page_content["notes"] = notes_slide.notes_text_frame.text

            if page_content["title"] or page_content["text_content"] or page_content["images"] or page_content["tables"]:
                pages_content[page_key] = page_content

        total_words = 0
        for page in pages_content.values():
            for text in page["text_content"]:
                total_words += len(text.split())
            if page["notes"]:
                total_words += len(page["notes"].split())

        result = {
            "type": "pptx",
            "metadata": metadata,
            "content": pages_content,
            "assessments": [],
            "answer_keys": [],
            "summary": {
                "total_pages": len(pages_content),
                "pages_with_images": len([p for p in pages_content.values() if p["images"]]),
                "pages_with_tables": len([p for p in pages_content.values() if p["tables"]]),
                "pages_with_notes": len([p for p in pages_content.values() if p["notes"]]),
                "total_text_blocks": sum(len(p["text_content"]) for p in pages_content.values()),
            },
            "total_word_count": total_words,
        }

        return result

    @staticmethod
    async def _parse_docx(file_content: bytes, filename: str) -> Dict[str, Any]:
        docx_file = io.BytesIO(file_content)
        document = DocxDocument(docx_file)

        metadata = {
            "filename": filename,
            "file_type": "DOCX",
            "parsed_at": datetime.utcnow().isoformat(),
            "total_paragraphs": len(document.paragraphs),
            "total_tables": len(document.tables),
            "total_images": len(document.inline_shapes),
            "title": document.core_properties.title or "",
            "author": document.core_properties.author or "",
            "subject": document.core_properties.subject or "",
            "created": document.core_properties.created.isoformat() if document.core_properties.created else "",
            "modified": document.core_properties.modified.isoformat() if document.core_properties.modified else "",
        }

        content_by_headings: Dict[str, Any] = {}
        current_heading: Optional[str] = None
        current_subheading: Optional[str] = None
        content_buffer: List[Dict[str, Any]] = []

        for paragraph in document.paragraphs:
            if not paragraph.text.strip():
                continue
            style_name = paragraph.style.name if paragraph.style else "Normal"
            is_heading = style_name.startswith("Heading")

            if is_heading:
                if current_heading and content_buffer:
                    if current_subheading:
                        if current_heading not in content_by_headings:
                            content_by_headings[current_heading] = {}
                        content_by_headings[current_heading][current_subheading] = content_buffer.copy()
                    else:
                        content_by_headings[current_heading] = content_buffer.copy()
                    content_buffer = []
                if style_name == "Heading 1":
                    current_heading = paragraph.text.strip()
                    current_subheading = None
                elif style_name == "Heading 2":
                    current_subheading = paragraph.text.strip()
                else:
                    if current_subheading:
                        current_subheading = f"{current_subheading} - {paragraph.text.strip()}"
                    else:
                        current_subheading = paragraph.text.strip()
            else:
                para_content = {
                    "text": paragraph.text.strip(),
                    "style": style_name,
                    "formatting": DocumentParser._extract_paragraph_formatting(paragraph),
                }
                content_buffer.append(para_content)

        if current_heading and content_buffer:
            if current_subheading:
                if current_heading not in content_by_headings:
                    content_by_headings[current_heading] = {}
                content_by_headings[current_heading][current_subheading] = content_buffer.copy()
            else:
                content_by_headings[current_heading] = content_buffer.copy()

        tables_data = []
        for table_num, table in enumerate(document.tables, 1):
            table_data = {
                "type": "table",
                "table_number": table_num,
                "rows": len(table.rows),
                "columns": len(table.columns) if table.rows else 0,
                "data": [],
            }
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data["data"].append(row_data)
            tables_data.append(table_data)

        images_data = []
        for i, shape in enumerate(document.inline_shapes, 1):
            image_data = {
                "type": "image",
                "image_number": i,
                "shape_type": str(shape.type),
                "width": shape.width.inches if shape.width else None,
                "height": shape.height.inches if shape.height else None,
            }
            images_data.append(image_data)

        total_words = 0
        for heading, sections in content_by_headings.items():
            if isinstance(sections, dict):
                for subheading, content in sections.items():
                    for item in content:
                        if isinstance(item, dict) and "text" in item:
                            total_words += len(item["text"].split())
            elif isinstance(sections, list):
                for item in sections:
                    if isinstance(item, dict) and "text" in item:
                        total_words += len(item["text"].split())
        for table in tables_data:
            for row in table["data"]:
                for cell in row:
                    total_words += len(cell.split())

        result = {
            "type": "docx",
            "metadata": metadata,
            "content": content_by_headings,
            "assessments": [],
            "answer_keys": [],
            "tables": tables_data,
            "images": images_data,
            "summary": {
                "total_headings": len(content_by_headings),
                "total_sections": sum(len(v) if isinstance(v, dict) else 1 for v in content_by_headings.values()),
                "total_paragraphs": sum(
                    len(v) if isinstance(v, list) else sum(len(sub) if isinstance(sub, list) else 0 for sub in v.values()) for v in content_by_headings.values()
                ),
                "total_tables": len(tables_data),
                "total_images": len(images_data),
            },
            "total_word_count": total_words,
        }

        return result

    @staticmethod
    def _extract_shape_content(shape) -> Optional[Dict[str, Any]]:
        try:
            shape_data: Dict[str, Any] = {"type": "unknown", "text": "", "position": {}, "size": {}}
            if hasattr(shape, "left") and hasattr(shape, "top"):
                shape_data["position"] = {"left": shape.left, "top": shape.top}
            if hasattr(shape, "width") and hasattr(shape, "height"):
                shape_data["size"] = {"width": shape.width, "height": shape.height}
            if getattr(shape, "has_text_frame", False):
                shape_data["type"] = "text"
                shape_data["text"] = getattr(shape, "text", "")
                if getattr(shape, "text_frame", None):
                    shape_data["paragraphs"] = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            shape_data["paragraphs"].append({
                                "text": para_text,
                                "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                                "level": paragraph.level,
                            })
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                shape_data["type"] = "image"
                shape_data["text"] = f"[Image: {shape.name if hasattr(shape, 'name') else 'Unknown'}]"
            elif getattr(shape, "has_table", False):
                shape_data["type"] = "table"
                table = shape.table
                shape_data["table_data"] = {"rows": len(table.rows), "columns": len(table.columns), "data": []}
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    shape_data["table_data"]["data"].append(row_data)
            else:
                shape_data["type"] = "shape"
                shape_data["text"] = shape.name if hasattr(shape, "name") else ""
            return shape_data if shape_data["text"] or shape_data["type"] in ["image", "table"] else None
        except Exception as e:
            logger.warning(f"Error extracting shape content: {str(e)}")
            return None

    @staticmethod
    def _extract_paragraph_formatting(paragraph) -> Dict[str, Any]:
        formatting: Dict[str, Any] = {"runs": []}
        for run in paragraph.runs:
            if run.text.strip():
                run_data = {
                    "text": run.text,
                    "bold": run.bold,
                    "italic": run.italic,
                    "underline": run.underline,
                    "font_size": str(run.font.size) if run.font.size else None,
                    "font_name": run.font.name if run.font.name else None,
                }
                formatting["runs"].append(run_data)
        return formatting

    @staticmethod
    def save_content_to_json(content: Dict[str, Any], file_path: str) -> str:
        with open(file_path, "w", encoding="utf-8") as f:
            json.dump(content, f, indent=2, ensure_ascii=False)

        diagnostics = content.get("diagnostics") or {}
        diag_lines: List[Dict[str, Any]] = []
        if diagnostics.get("front_matter_pages"):
            diag_lines.append({"type": "front_matter_pages", "data": diagnostics["front_matter_pages"]})
        if diagnostics.get("mismatch_pages"):
            diag_lines.append({"type": "mismatch_pages", "data": diagnostics["mismatch_pages"]})
        if diagnostics.get("auto_repairs"):
            diag_lines.append({"type": "auto_repairs", "data": diagnostics["auto_repairs"]})
        if diagnostics.get("heading_candidates_rejected"):
            diag_lines.append({
                "type": "heading_candidates_rejected",
                "data": diagnostics["heading_candidates_rejected"],
            })
        if diagnostics.get("out_of_range_headings"):
            diag_lines.append({"type": "out_of_range_headings", "data": diagnostics["out_of_range_headings"]})
        if diagnostics.get("figure_caption_orphans"):
            diag_lines.append({"type": "figure_caption_orphans", "data": diagnostics["figure_caption_orphans"]})
        if diagnostics.get("ocr_confidence_histogram"):
            diag_lines.append({"type": "ocr_confidence_histogram", "data": diagnostics["ocr_confidence_histogram"]})
        if diagnostics.get("no_toc_lecture_mode_used"):
            diag_lines.append({"type": "no_toc_lecture_mode_used", "data": diagnostics["no_toc_lecture_mode_used"]})
        if diagnostics.get("synthesized_ranges"):
            diag_lines.append({"type": "synthesized_ranges", "data": diagnostics["synthesized_ranges"]})
        if diagnostics.get("synonym_mappings"):
            diag_lines.append({"type": "synonym_mappings", "data": diagnostics["synonym_mappings"]})
        if diagnostics.get("unmapped_headings"):
            diag_lines.append({"type": "unmapped_headings", "data": diagnostics["unmapped_headings"]})
        if diagnostics.get("answer_key_pages_found"):
            diag_lines.append({"type": "answer_key_pages_found", "data": diagnostics["answer_key_pages_found"]})
        if diagnostics.get("assessment_blocks_found"):
            diag_lines.append({"type": "assessment_blocks_found", "data": diagnostics["assessment_blocks_found"]})
        if diagnostics.get("orphan_answers"):
            diag_lines.append({"type": "orphan_answers", "data": diagnostics["orphan_answers"]})
        if diagnostics.get("unmatched_questions"):
            diag_lines.append({"type": "unmatched_questions", "data": diagnostics["unmatched_questions"]})

        if diag_lines:
            diag_path = Path(file_path).with_suffix(".diagnostics.jsonl")
            with open(diag_path, "w", encoding="utf-8") as diag_file:
                for line in diag_lines:
                    diag_file.write(json.dumps(line, ensure_ascii=False) + "\n")

        return file_path


if __name__ == "__main__":
    print("DocumentParser service loaded successfully!")
    print("Available document types:")
    for doc_type in DocumentType:
        print(f"  • {doc_type.value.upper()}")
    print("\nDocumentParser class is ready to use.")
